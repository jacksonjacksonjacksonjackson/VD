"""
powerpoint_export.py

Template-driven PowerPoint export for the Fleet Electrification Analyzer.
Loads the bundled template PPTX (assets/template_default.pptx), modifies slides
in-place with fleet data, and saves to the requested output path.

Architecture: Template-Modify
- Static slides (CARB overview, ACF exemptions, incentives) are kept as-is.
- Token slides (cover, contact) have text placeholders updated with profile fields.
- Chart slides have their embedded chart replaced with live fleet data.
- Key Findings slide gets auto-generated bullets from the processed fleet.
- Unchecked slides are deleted in reverse-index order.
- Optional extra slides are appended using the template's slide master layout.
"""

import os
import json
import logging
import datetime
import subprocess
import shutil
from collections import defaultdict
from typing import Dict, List, Any, Optional
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from data.models import Fleet, FleetVehicle, PresentationProfile
from settings import (
    APP_NAME, APP_VERSION,
    PRIMARY_HEX_1, PRIMARY_HEX_2, PRIMARY_HEX_3,
    SECONDARY_HEX_1, SECONDARY_HEX_2,
    DEFAULT_TEMPLATE_PATH, TEMPLATE_SLIDE_IDS, DEFAULT_SLIDE_IDS,
    EXPORT_DIR,
)
from powerpoint_charts import (
    SlideConfiguration,
    add_fleet_composition_chart,
    add_acf_category_composition_chart,
    add_department_summary_chart,
    add_facility_summary_chart,
    add_electrification_timeline_by_weight_chart,
    add_tco_comparison_chart,
    add_payback_timeline_chart,
    add_scenario_comparison_chart,
    add_age_distribution_chart,
)

logger = logging.getLogger(__name__)

###############################################################################
# Colour helpers
###############################################################################

def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# Brand colours derived from settings
_C1 = _hex_to_rgb(PRIMARY_HEX_1)   # charcoal
_C3 = _hex_to_rgb(PRIMARY_HEX_3)   # reseda green
_S1 = _hex_to_rgb(SECONDARY_HEX_1) # deep orange
_GREY = _hex_to_rgb(SECONDARY_HEX_2)

###############################################################################
# Low-level slide manipulation helpers
###############################################################################

def _remove_chart_shapes(slide) -> Optional[dict]:
    """Remove all chart shapes from a slide. Returns the last removed chart's bounds."""
    bounds = None
    for shape in list(slide.shapes):
        if shape.has_chart:
            bounds = {
                "left": shape.left, "top": shape.top,
                "width": shape.width, "height": shape.height,
            }
            shape._element.getparent().remove(shape._element)
    return bounds


def _delete_slide(prs: Presentation, slide_index: int) -> None:
    """Delete a slide by (0-based) index."""
    xml_slides = prs.slides._sldIdLst
    if slide_index >= len(xml_slides):
        return
    elem = xml_slides[slide_index]
    ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rId = elem.get(f"{{{ns}}}id")
    if rId:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
    xml_slides.remove(elem)


def _reorder_slides(prs: Presentation, ordered_ids: List[str]) -> None:
    """
    Reorder slides so they match ordered_ids (subset of TEMPLATE_SLIDE_IDS).
    Slides not in ordered_ids are silently skipped.
    """
    xml_slides = prs.slides._sldIdLst
    items = list(xml_slides)
    # Build id→xml-element map.
    # After deletions, the remaining slides are in template order but with gaps.
    # Map by position against the filtered (remaining) template ID list, not the
    # full TEMPLATE_SLIDE_IDS list — otherwise indices shift after any deletion.
    ordered_set = set(ordered_ids)
    remaining_ids = [sid for sid in TEMPLATE_SLIDE_IDS if sid in ordered_set]
    id_to_elem = {
        remaining_ids[i]: items[i]
        for i in range(min(len(remaining_ids), len(items)))
    }
    # Remove all
    for item in items:
        xml_slides.remove(item)
    # Re-add in requested order
    for slide_id in ordered_ids:
        if slide_id in id_to_elem:
            xml_slides.append(id_to_elem[slide_id])


def _get_layout(prs: Presentation, name: str):
    """Get a slide layout by name from the first slide master."""
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_masters[0].slide_layouts[1]  # fallback: "Title and Content"


def _find_placeholder(slide, idx: int):
    """Return placeholder by idx, or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_text_frame_lines(tf, lines: List[str], keep_format: bool = True) -> None:
    """
    Replace the content of a text frame with lines.
    Optionally copies the first run's font onto new runs.
    """
    # Capture reference font from first run
    ref_size = None
    ref_bold = None
    ref_color = None
    if keep_format:
        for para in tf.paragraphs:
            for run in para.runs:
                ref_size = run.font.size
                ref_bold = run.font.bold
                try:
                    ref_color = run.font.color.rgb
                except Exception:
                    pass
                break
            if ref_size is not None:
                break

    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
        para.text = line
        if keep_format and line:
            run = para.runs[0] if para.runs else para.add_run()
            if ref_size:
                run.font.size = ref_size
            if ref_bold is not None:
                run.font.bold = ref_bold


def _replace_in_text_frame(tf, old: str, new: str) -> None:
    """Replace all occurrences of old with new across all runs in a text frame."""
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def _replace_tokens_in_slide(slide, token_map: dict) -> None:
    """Token-replace across all text frames in a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for old, new in token_map.items():
            _replace_in_text_frame(shape.text_frame, old, new)


###############################################################################
# Chart style helpers
###############################################################################

def _style_chart(chart, legend_pos=None) -> None:
    """Apply clean chart styling: no title, clean legend, minimal gridlines."""
    chart.has_title = False
    if legend_pos is not None:
        chart.has_legend = True
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False


def _hex_series_colors(chart, colors: List[str]) -> None:
    """Apply a list of hex colour strings to chart series, cycling if needed."""
    if not PPTX_AVAILABLE:
        return
    series_list = list(chart.series)
    for i, series in enumerate(series_list):
        colour = _hex_to_rgb(colors[i % len(colors)])
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colour
        try:
            series.format.line.color.rgb = colour
        except Exception:
            pass


# Default colour cycle for charts
_CHART_COLORS = [
    PRIMARY_HEX_3,   # reseda green
    PRIMARY_HEX_1,   # charcoal
    SECONDARY_HEX_1, # deep orange
    "#7B9E87",       # sage
    "#A8C5DA",       # sky blue
    "#D4A373",       # warm sand
]

###############################################################################
# Key Findings generator
###############################################################################

def _generate_key_findings(vehicles: List[FleetVehicle], profile: PresentationProfile) -> List[str]:
    """Auto-generate narrative key finding bullets from fleet data.

    Produces client-facing language (not technical) with conditional bullets
    that appear only when data supports them. Editable in PowerPoint after export.
    """
    client = profile.client_name or "the fleet"
    n = len(vehicles)
    findings = []

    # 1. Fleet inventory overview
    findings.append(
        f"{client} provided a fleet inventory of {n} vehicle{'s' if n != 1 else ''}."
    )

    # 2. Department / facility context (only if data present)
    depts = {(v.department or "").strip() for v in vehicles if (v.department or "").strip()}
    locs = {(getattr(v, "location", "") or "").strip() for v in vehicles
            if (getattr(v, "location", "") or "").strip()}
    if depts and locs:
        findings.append(
            f"The fleet operates across {len(depts)} department{'s' if len(depts) != 1 else ''} "
            f"and {len(locs)} domicile facilit{'ies' if len(locs) != 1 else 'y'}."
        )
    elif depts:
        findings.append(
            f"The fleet operates across {len(depts)} department{'s' if len(depts) != 1 else ''}."
        )

    # Build ACF counts
    acf_counts: Dict[str, int] = defaultdict(int)
    for v in vehicles:
        code = v.custom_fields.get("_acf_code", "")
        acf_counts[code] += 1

    # 3. ZEV vehicles (only if present)
    zev_n = acf_counts.get("ZEV", 0)
    if zev_n > 0:
        findings.append(
            f"{zev_n} vehicle{'s are' if zev_n != 1 else ' is'} already zero-emission "
            f"and classified as ZEV under CARB ACF."
        )

    # 3b. EV-replaceable % (Cat B mandate-subject)
    b_count_pre = acf_counts.get("B", 0)
    if b_count_pre > 0 and n > 0:
        pct = round(100 * b_count_pre / n)
        findings.append(
            f"{b_count_pre} of {n} vehicles ({pct}%) are classified as mandate-subject "
            f"(Category B — medium and heavy duty) and subject to CARB ACF replacement timelines."
        )

    # 4. Exemptions (Cat C and D)
    c_count = acf_counts.get("C", 0)
    d_count = acf_counts.get("D", 0)
    a_count = acf_counts.get("A", 0)
    exempt_parts = []
    if d_count:
        exempt_parts.append(
            f"{d_count} emergency vehicle{'s' if d_count != 1 else ''} (Category D)"
        )
    if c_count:
        exempt_parts.append(
            f"{c_count} body-type exempt vehicle{'s' if c_count != 1 else ''} (Category C)"
        )
    if a_count:
        exempt_parts.append(
            f"{a_count} light-duty vehicle{'s' if a_count != 1 else ''} under 8,500 lbs GVWR (Category A)"
        )
    if exempt_parts:
        findings.append(
            f"The fleet includes {', and '.join(exempt_parts)} that are exempt "
            f"from CARB Advanced Clean Fleets mandate-year requirements."
        )

    # 5. ACF-B urgency: earliest mandate deadline year
    b_count = acf_counts.get("B", 0)
    if b_count > 0:
        b_years = []
        for v in vehicles:
            if v.custom_fields.get("_acf_code") == "B":
                try:
                    b_years.append(int(v.custom_fields.get("Proposed EV Year", "")))
                except (ValueError, TypeError):
                    pass
        if b_years:
            findings.append(
                f"The earliest mandate-subject replacement is projected for {min(b_years)}, "
                f"with procurement planning recommended to begin immediately."
            )

    # 6. Timeline achievability (all vehicles)
    ev_years = []
    for v in vehicles:
        try:
            yr = int(v.custom_fields.get("Proposed EV Year", ""))
            if 2024 <= yr <= 2060:
                ev_years.append(yr)
        except (ValueError, TypeError):
            pass

    if ev_years and b_count > 0:
        findings.append(
            f"Adhering to the current replacement schedule, 100% electrification of "
            f"mandate-subject vehicles is achievable by {max(ev_years)}."
        )

    return findings

###############################################################################
# Per-year emissions calculator
###############################################################################

def _calculate_yearly_emissions(vehicles: List[FleetVehicle]) -> Dict[int, float]:
    """
    For each calendar year, calculate total remaining ICE fleet emissions (metric tons CO2e).
    A vehicle contributes ICE emissions until its Proposed EV Year.
    Returns {year: metric_tons_co2} from current year to the last EV replacement year.
    """
    current_year = datetime.datetime.now().year
    ev_years = []
    for v in vehicles:
        try:
            yr = int(v.custom_fields.get("Proposed EV Year", ""))
            if 2024 <= yr <= 2060:
                ev_years.append(yr)
        except (ValueError, TypeError):
            pass

    if not ev_years:
        return {}

    max_year = max(ev_years)
    result = {}
    for year in range(current_year, max_year + 1):
        total = 0.0
        for v in vehicles:
            try:
                ev_yr = int(v.custom_fields.get("Proposed EV Year", ""))
            except (ValueError, TypeError):
                ev_yr = max_year + 1  # Never replaced
            if year < ev_yr:
                co2_g_mi = v.fuel_economy.co2_primary or 0.0
                miles = v.annual_mileage or 12000
                total += (co2_g_mi * miles) / 1_000_000  # g → metric tons
        result[year] = round(total, 2)
    return result

###############################################################################
# Slide-specific updaters
###############################################################################

def _update_cover_slide(slide, profile: PresentationProfile) -> None:
    """Update cover slide subtitle with client name, date, presenter."""
    # Build subtitle lines
    ptype = profile.presentation_type or "Kickoff"
    subtitle_lines = [
        f"{profile.client_name} {ptype}" if profile.client_name else ptype,
        profile.meeting_date or "",
        "",
        "Presented by",
        profile.presenter_name or "",
    ]
    # Placeholder idx 0 = title ("Plan Your Fleet") – leave as-is
    # Placeholder idx 1 = subtitle – rebuild
    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame:
        _set_text_frame_lines(ph.text_frame, subtitle_lines)


def _update_key_findings_slide(slide, vehicles: List[FleetVehicle], profile: PresentationProfile) -> None:
    """Replace the Key Findings bullet list with auto-generated findings."""
    bullets = _generate_key_findings(vehicles, profile)
    # Placeholder idx 1 = content body
    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame:
        _set_text_frame_lines(ph.text_frame, bullets)


def _update_text_list_slide(slide, items: List[str]) -> None:
    """Update a content-placeholder slide with a list of bullet strings."""
    if not items:
        return
    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame:
        _set_text_frame_lines(ph.text_frame, items)


def _update_contact_slide(slide, profile: PresentationProfile) -> None:
    """
    Update the contact slide with partner and presenter info.
    The template has:
      Placeholder idx 1 (subtitle) → partner 1
      A free TextBox shape        → partner 2 / presenter
    """
    # Build partner 1 text
    p1_lines = [l for l in [
        profile.partner_1_name,
        profile.partner_1_title + (f", {profile.partner_1_org}" if profile.partner_1_org else ""),
        profile.partner_1_email,
    ] if l.strip()]

    # Build partner 2 / presenter text
    p2_lines = [l for l in [
        profile.partner_2_name or profile.presenter_name,
        (profile.partner_2_title or profile.presenter_title)
            + (f", {profile.partner_2_org or profile.presenter_company}"
               if (profile.partner_2_org or profile.presenter_company) else ""),
        profile.partner_2_email,
    ] if l.strip()]

    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame and p1_lines:
        _set_text_frame_lines(ph.text_frame, p1_lines)

    # The free TextBox is typically the last shape with a text frame that isn't a placeholder
    text_boxes = [s for s in slide.shapes if s.has_text_frame
                  and not s.is_placeholder]
    if text_boxes and p2_lines:
        _set_text_frame_lines(text_boxes[-1].text_frame, p2_lines)

###############################################################################
# Chart generators (template-slot replacements)
###############################################################################

def _add_acf_electrification_chart(
    slide,
    vehicles: List[FleetVehicle],
    year_map: Optional[Dict[str, str]] = None,
    end_year: int = 2040,
) -> bool:
    """Stacked column chart: vehicles replaced per year, grouped by ACF category.

    Matches the style of the user-built deck (Slide 8).
    All 4 ACF categories are always present (even if zero in some years) so the
    legend is consistent across scenario slides.

    Args:
        year_map: Optional {vin: year_str} override for scenario slides. If None,
                  reads "Proposed EV Year" from custom_fields (current plan).
        end_year: Chart X-axis ceiling year (default 2040). Extended dynamically
                  if any assignments exceed this value.
    """
    if not PPTX_AVAILABLE:
        return False

    current_year = datetime.datetime.now().year

    # Plain-English labels matching the user-built deck
    ACF_LABELS = {
        "A": "Light Duty (Excluded)",
        "B": "Medium or Heavy Duty",
        "C": "Medium or Heavy Duty Potentially Exempted",
        "D": "Emergency Vehicles (Excluded)",
    }
    # Display order: B first (mandate-subject), then exempt categories
    ACF_ORDER = ["B", "A", "C", "D"]
    # Colors: B=charcoal (prominent), A=sky blue, C=sage, D=orange
    ACF_COLORS = [PRIMARY_HEX_1, "#A8C5DA", "#7B9E87", SECONDARY_HEX_1]

    # Collect year × ACF counts
    year_acf: Dict[int, Dict[str, int]] = defaultdict(lambda: defaultdict(int))
    present_codes: set = set()

    for v in vehicles:
        code = v.custom_fields.get("_acf_code", "")
        if code not in ACF_LABELS:
            continue  # skip ZEV, unclassified — they don't get replacement years
        # Get year from override map or current plan
        if year_map is not None:
            year_str = year_map.get(v.vin, "")
        else:
            year_str = v.custom_fields.get("Proposed EV Year", "")
        try:
            yr = int(year_str)
            if current_year <= yr <= 2060:
                year_acf[yr][code] += 1
                present_codes.add(code)
        except (ValueError, TypeError):
            pass

    if not year_acf:
        return False

    # Cap at end_year; extend only if actual assignments exceed it
    chart_end = max(end_year, max(year_acf.keys()))
    start_year = min(current_year, 2026)
    years = list(range(start_year, chart_end + 1))

    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in years]

    # Always add all 4 ACF categories so legends are consistent across scenario slides
    for code in ACF_ORDER:
        values = [year_acf[y].get(code, 0) for y in years]
        chart_data.add_series(ACF_LABELS[code], values)

    bounds = _remove_chart_shapes(slide) or {
        "left": Inches(2.16), "top": Inches(1.85),
        "width": Inches(10.11), "height": Inches(4.89),
    }

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        bounds["left"], bounds["top"], bounds["width"], bounds["height"],
        chart_data,
    ).chart

    _style_chart(chart, legend_pos=XL_LEGEND_POSITION.BOTTOM)
    _hex_series_colors(chart, ACF_COLORS)

    # Data labels: show vehicle count per bar segment; suppress zeros via numFmt
    try:
        for series in chart.series:
            series.data_labels.show_value = True
            series.data_labels.number_format = "#,##0;-#,##0;;"  # empty for zero
    except Exception:
        pass

    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = _hex_to_rgb("#E0E0E0")
    except Exception:
        pass

    return True


def _calculate_baseline_emissions(vehicles: List[FleetVehicle]) -> float:
    """Total annual fleet emissions if NO electrification occurs (all stay ICE)."""
    total = 0.0
    for v in vehicles:
        co2_g_mi = v.fuel_economy.co2_primary or 0.0
        if co2_g_mi <= 0:
            mpg = v.fuel_economy.combined_mpg or 0.0
            co2_g_mi = 8900 / mpg if mpg > 0 else 0.0
        miles = v.annual_mileage or 12000
        total += (co2_g_mi * miles) / 1_000_000
    return round(total, 2)


def _add_ghg_emissions_chart(slide, vehicles: List[FleetVehicle]) -> bool:
    """Line chart: 3-series GHG emissions comparison matching the user-built deck.

    Series:
      1. Baseline (No Electrification) — flat line at total annual fleet CO₂e
      2. Medium or Heavy Duty Only — only Cat B emissions declining as B vehicles electrify
      3. Whole Fleet Electrification — total fleet emissions declining (current plan)

    X-axis: 2026–2040 (extended if vehicles scheduled beyond 2040).
    """
    if not PPTX_AVAILABLE:
        return False

    current_year = datetime.datetime.now().year
    start_year = min(current_year, 2026)
    end_year_default = 2040

    # Whole fleet: all vehicles declining as they get EV years
    whole_fleet_yearly = _calculate_yearly_emissions(vehicles)
    if not whole_fleet_yearly:
        return False

    chart_end = max(end_year_default, max(whole_fleet_yearly.keys()))
    years = list(range(start_year, chart_end + 1))

    # Baseline: flat — total fleet CO₂e if nothing is electrified
    baseline = _calculate_baseline_emissions(vehicles)
    baseline_values = [baseline] * len(years)

    # Medium/Heavy Duty Only: only Cat B vehicles electrified; A/C/D stay at ICE
    b_vehicles = [v for v in vehicles if v.custom_fields.get("_acf_code") == "B"]
    b_yearly = _calculate_yearly_emissions(b_vehicles) if b_vehicles else {}
    # For years before current_year (chart can show 2026 even if year is 2027),
    # fall back to the full static B-fleet baseline (no vehicle replaced yet).
    b_baseline = sum(
        (v.fuel_economy.co2_primary or 0) * (v.annual_mileage or 12000) / 1_000_000
        for v in b_vehicles
    )
    b_values = [round(b_yearly.get(y, b_baseline), 2) for y in years]

    whole_values = [whole_fleet_yearly.get(y, 0.0) for y in years]

    chart_data = CategoryChartData()
    chart_data.categories = [str(y) for y in years]
    chart_data.add_series("Baseline (No Electrification)", baseline_values)
    if b_vehicles:
        chart_data.add_series("Medium or Heavy Duty Only", b_values)
    chart_data.add_series("Whole Fleet Electrification", whole_values)

    bounds = _remove_chart_shapes(slide) or {
        "left": Inches(2.25), "top": Inches(1.73),
        "width": Inches(9.89), "height": Inches(5.20),
    }

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        bounds["left"], bounds["top"], bounds["width"], bounds["height"],
        chart_data,
    ).chart

    _style_chart(chart, legend_pos=XL_LEGEND_POSITION.BOTTOM)
    # Baseline=gray, M/H Only=charcoal, Whole Fleet=green
    emission_colors = [SECONDARY_HEX_2, PRIMARY_HEX_1, PRIMARY_HEX_3]
    if not b_vehicles:
        emission_colors = [SECONDARY_HEX_2, PRIMARY_HEX_3]
    _hex_series_colors(chart, emission_colors)

    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = _hex_to_rgb("#E0E0E0")
    except Exception:
        pass

    try:
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.text = "Metric Tons CO\u2082e"
    except Exception:
        pass

    return True


def _add_tco_appendix_chart(slide, vehicles: List[FleetVehicle]) -> bool:
    """Clustered column: annual marginal EV vs ICE total cost of ownership."""
    if not PPTX_AVAILABLE:
        return False

    bounds = _remove_chart_shapes(slide) or {
        "left": Inches(2.31), "top": Inches(1.94),
        "width": Inches(9.78), "height": Inches(4.23),
    }
    result = add_tco_comparison_chart(
        slide, vehicles,
        left=bounds["left"] / 914400,
        top=bounds["top"] / 914400,
        width=bounds["width"] / 914400,
        height=bounds["height"] / 914400,
    )
    if result:
        # Style the newly added chart
        for shape in slide.shapes:
            if shape.has_chart:
                _style_chart(shape.chart, legend_pos=XL_LEGEND_POSITION.BOTTOM)
                _hex_series_colors(shape.chart, [PRIMARY_HEX_3, SECONDARY_HEX_1])
    return result


###############################################################################
# Optional extra slide builders
###############################################################################

def _append_optional_slide(prs: Presentation, slide_id: str,
                           vehicles: List[FleetVehicle],
                           profile: PresentationProfile,
                           scenario_results: Optional[list] = None) -> bool:
    """Add an optional extra slide to the presentation.

    Returns True if a slide was successfully added, False otherwise.
    Slide insertion position is handled by the caller (export_presentation).
    """
    if not PPTX_AVAILABLE:
        return False

    # ── Phase 24 new slide types ──────────────────────────────────────────────
    if slide_id == "acf_composition":
        return _create_acf_composition_slide(prs, vehicles)

    if slide_id in ("timeline_moderate", "timeline_aggressive",
                    "timeline_conservative", "timeline_current_plan"):
        return _create_scenario_timeline_slide(prs, vehicles, slide_id)

    if slide_id == "invalid_vin":
        return _create_invalid_vin_slide(prs, vehicles)

    if slide_id == "department_summary":
        return _create_department_summary_slide(prs, vehicles)

    if slide_id == "facility_summary":
        return _create_facility_summary_slide(prs, vehicles)

    # ── Existing optional slide types ─────────────────────────────────────────
    layout = _get_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)

    if slide_id in ("fleet_composition", "fleet_composition_body_type"):
        _set_placeholder_title(slide, "Fleet Composition by Body Type")
        return add_fleet_composition_chart(slide, vehicles, left=2.0, top=1.8, width=9.0, height=5.0)

    elif slide_id == "age_analysis":
        _set_placeholder_title(slide, "Fleet Age Distribution")
        return add_age_distribution_chart(slide, vehicles, left=2.0, top=1.8, width=9.0, height=5.0)

    elif slide_id == "scenario_comparison":
        if scenario_results:
            _set_placeholder_title(slide, "Electrification Scenario Comparison")
            return add_scenario_comparison_chart(slide, scenario_results, left=2.0, top=1.8, width=9.0, height=5.0)
        # No data — remove the blank slide we just added
        _delete_slide(prs, len(prs.slides) - 1)
        return False

    elif slide_id == "replacement_table":
        _set_placeholder_title(slide, "Priority Vehicle Replacement Schedule")
        _add_replacement_table_to_slide(slide, vehicles)
        return True

    elif slide_id == "data_quality":
        _set_placeholder_title(slide, "Data Quality & Completeness")
        _add_data_quality_text(slide, vehicles)
        return True

    # Unknown ID — remove the blank slide
    _delete_slide(prs, len(prs.slides) - 1)
    return False


def _set_placeholder_title(slide, title_text: str) -> None:
    ph = _find_placeholder(slide, 0)
    if ph and ph.has_text_frame:
        ph.text_frame.text = title_text


def _add_replacement_table_to_slide(slide, vehicles: List[FleetVehicle]) -> None:
    """Add a simple text summary of the top-priority vehicles."""
    # Sort by proposed EV year (earliest first)
    sorted_veh = sorted(vehicles, key=lambda v: (
        int(v.custom_fields.get("Proposed EV Year", "9999"))
        if v.custom_fields.get("Proposed EV Year", "").isdigit() else 9999
    ))[:12]

    lines = ["Asset ID | Make | Model | Year | ACF | Target EV Year"]
    lines.append("—" * 55)
    for v in sorted_veh:
        asset = v.custom_fields.get("Asset ID", v.vin[-6:])
        acf = v.custom_fields.get("ACF Category", "—")
        ev_yr = v.custom_fields.get("Proposed EV Year", "—")
        lines.append(f"{asset:12} | {v.vehicle_id.make:8} | {v.vehicle_id.model:12} | "
                     f"{v.vehicle_id.year:4} | {acf:5} | {ev_yr}")

    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame:
        _set_text_frame_lines(ph.text_frame, lines)


def _add_data_quality_text(slide, vehicles: List[FleetVehicle]) -> None:
    """Add data quality summary bullets."""
    n = len(vehicles)
    if n == 0:
        return
    vin_ok = sum(1 for v in vehicles if v.vehicle_id.make)
    mpg_ok = sum(1 for v in vehicles if v.fuel_economy.combined_mpg > 0)
    acf_ok = sum(1 for v in vehicles if v.custom_fields.get("_acf_code"))

    lines = [
        f"Fleet Size: {n} vehicles analysed",
        f"VIN Decode Coverage: {vin_ok/n:.0%} ({vin_ok}/{n})",
        f"MPG Data Coverage: {mpg_ok/n:.0%} ({mpg_ok}/{n})",
        f"ACF Classification: {acf_ok/n:.0%} ({acf_ok}/{n})",
    ]
    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame:
        _set_text_frame_lines(ph.text_frame, lines)

###############################################################################
# Slide insertion utility (Phase 24)
###############################################################################

def _move_slide_to_index(prs: Presentation, from_idx: int, to_idx: int) -> None:
    """Move the slide at from_idx to to_idx in the presentation."""
    xml_slides = prs.slides._sldIdLst
    items = list(xml_slides)
    if from_idx == to_idx or from_idx >= len(items):
        return
    elem = items[from_idx]
    xml_slides.remove(elem)
    # After removal, indices shift; correct to_idx if it was after from_idx
    adj = to_idx if to_idx <= from_idx else to_idx - 1
    xml_slides.insert(adj, elem)


###############################################################################
# New optional slide builders (Phase 24)
###############################################################################

def _create_acf_composition_slide(prs: Presentation, vehicles: List[FleetVehicle]) -> bool:
    """Fleet Composition by ACF Category — pie chart optional slide."""
    layout = _get_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    _set_placeholder_title(slide, "Fleet Composition by ACF Category")
    return add_acf_category_composition_chart(
        slide, vehicles, left=1.5, top=1.8, width=10.0, height=5.2
    )


def _create_scenario_timeline_slide(
    prs: Presentation,
    vehicles: List[FleetVehicle],
    slide_id: str,
) -> bool:
    """Electrification Timeline for a specific scenario — stacked column optional slide.

    slide_id options: timeline_moderate, timeline_aggressive, timeline_conservative,
                      timeline_current_plan
    """
    from analysis.scenarios import get_scenario_year_assignments

    SCENARIO_META = {
        "timeline_moderate":      ("Electrification Timeline — Moderate (2035 Target)",     "moderate",     2035),
        "timeline_aggressive":    ("Electrification Timeline — Aggressive (2030 Target)",   "aggressive",   2030),
        "timeline_conservative":  ("Electrification Timeline — Conservative (2040 Target)", "conservative", 2040),
        "timeline_current_plan":  ("Electrification Timeline — Current Plan",               None,           2040),
    }
    meta = SCENARIO_META.get(slide_id)
    if not meta:
        return False

    title, scenario_key, end_yr = meta
    layout = _get_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    _set_placeholder_title(slide, title)

    if scenario_key is None:
        # Current plan: read from custom_fields directly
        year_map = None
    else:
        year_map = get_scenario_year_assignments(vehicles, scenario_key)

    return _add_acf_electrification_chart(slide, vehicles, year_map=year_map, end_year=end_yr)


def _create_invalid_vin_slide(prs: Presentation, vehicles: List[FleetVehicle]) -> bool:
    """Vehicle Data Assumptions — lists vehicles where VIN decoding failed."""
    invalid = [v for v in vehicles if not v.vehicle_id.make]
    if not invalid:
        return False

    layout = _get_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    _set_placeholder_title(slide, "Vehicle Data Assumptions")

    lines = [
        "The following vehicles could not be fully decoded from VIN. "
        "Where possible, make/model/year were assumed from supporting data.",
        "",
    ]
    for v in invalid[:20]:  # cap at 20 rows
        asset = v.custom_fields.get("Asset ID") or v.vin[-6:] if v.vin else "—"
        make = v.vehicle_id.make or "Unknown Make"
        model = v.vehicle_id.model or "Unknown Model"
        year = str(v.vehicle_id.year) if v.vehicle_id.year else "—"
        ev_yr = v.custom_fields.get("Proposed EV Year", "—")
        lines.append(f"{asset}  |  {year} {make} {model}  |  Proposed EV Year: {ev_yr}")

    if len(invalid) > 20:
        lines.append(f"… and {len(invalid) - 20} additional vehicles (see exported CSV for full list).")

    ph = _find_placeholder(slide, 1)
    if ph and ph.has_text_frame:
        _set_text_frame_lines(ph.text_frame, lines)
    return True


def _create_department_summary_slide(prs: Presentation, vehicles: List[FleetVehicle]) -> bool:
    """Department Summary — horizontal bar chart, conditional on dept data."""
    has_dept = any((v.department or "").strip() for v in vehicles)
    if not has_dept:
        return False

    layout = _get_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    _set_placeholder_title(slide, "Department Summary")
    return add_department_summary_chart(slide, vehicles, left=1.5, top=1.8, width=10.0, height=5.2)


def _create_facility_summary_slide(prs: Presentation, vehicles: List[FleetVehicle]) -> bool:
    """Domicile Facility Summary — horizontal bar chart, conditional on location data."""
    has_loc = any((getattr(v, "location", "") or "").strip() for v in vehicles)
    if not has_loc:
        return False

    layout = _get_layout(prs, "Title and Content")
    slide = prs.slides.add_slide(layout)
    _set_placeholder_title(slide, "Domicile Facility Summary")
    return add_facility_summary_chart(slide, vehicles, left=1.5, top=1.8, width=10.0, height=5.2)


###############################################################################
# Main export function
###############################################################################

def export_presentation(
    fleet_data,                          # Fleet object or List[FleetVehicle]
    profile: Optional[PresentationProfile] = None,
    out_path: Optional[str] = None,
    template_path: Optional[str] = None,
    scenario_results: Optional[list] = None,
) -> str:
    """
    Generate a client-facing .pptx by modifying the template in-place with fleet data.

    Args:
        fleet_data:       Fleet object or list of FleetVehicle.
        profile:          PresentationProfile with client/consultant info.
        out_path:         Output file path; auto-generated if None.
        template_path:    Path to .pptx template; falls back to DEFAULT_TEMPLATE_PATH.
        scenario_results: Optional list of scenario result dicts for the optional
                          scenario_comparison slide.

    Returns:
        Absolute path to the saved .pptx file.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required. Run: pip install python-pptx")

    # --- Normalise inputs ---
    if profile is None:
        profile = PresentationProfile()
        profile.included_slides = list(DEFAULT_SLIDE_IDS)

    if isinstance(fleet_data, Fleet):
        vehicles = fleet_data.vehicles
    elif isinstance(fleet_data, list):
        vehicles = fleet_data
    else:
        vehicles = []

    tpl = template_path or profile.template_path or DEFAULT_TEMPLATE_PATH
    if not os.path.isfile(tpl):
        raise FileNotFoundError(f"Template not found: {tpl}")

    logger.info("[export_presentation] Loading template: %s", tpl)
    prs = Presentation(tpl)

    # --- Determine slide order ---
    included = profile.included_slides if profile.included_slides else list(DEFAULT_SLIDE_IDS)
    # Keep only IDs that exist in the template
    included = [sid for sid in included if sid in TEMPLATE_SLIDE_IDS]

    # --- Modify template slides in-place ---
    # Build a map: slide_id → current slide object (before any deletion)
    slide_map = {sid: prs.slides[i] for i, sid in enumerate(TEMPLATE_SLIDE_IDS)
                 if i < len(prs.slides)}

    # Cover
    if "cover" in slide_map:
        _update_cover_slide(slide_map["cover"], profile)

    # Key Findings
    if "key_findings" in slide_map and vehicles:
        _update_key_findings_slide(slide_map["key_findings"], vehicles, profile)

    # Electrification Timeline Chart (slide 7)
    if "timeline_chart" in slide_map and vehicles:
        _remove_chart_shapes(slide_map["timeline_chart"])
        _add_acf_electrification_chart(slide_map["timeline_chart"], vehicles)

    # GHG Emissions Chart (slide 8)
    if "emissions_chart" in slide_map and vehicles:
        _remove_chart_shapes(slide_map["emissions_chart"])
        _add_ghg_emissions_chart(slide_map["emissions_chart"], vehicles)

    # Data Needs
    if "data_needs" in slide_map and profile.data_needs_items:
        _update_text_list_slide(slide_map["data_needs"], profile.data_needs_items)

    # Next Steps
    if "next_steps" in slide_map and profile.next_steps_items:
        _update_text_list_slide(slide_map["next_steps"], profile.next_steps_items)

    # Contact
    if "contact" in slide_map:
        _update_contact_slide(slide_map["contact"], profile)

    # Charging Infra Costs (slide 14) – keep static; dept data not tracked yet
    # Annual Marginal EV TCO (slide 15)
    if "tco_chart" in slide_map and vehicles:
        _remove_chart_shapes(slide_map["tco_chart"])
        _add_tco_appendix_chart(slide_map["tco_chart"], vehicles)

    # --- Delete unchecked template slides (reverse order) ---
    excluded_template = [sid for sid in TEMPLATE_SLIDE_IDS if sid not in included]
    # Get indices in the CURRENT template ordering (before any reorder)
    indices_to_delete = sorted(
        [TEMPLATE_SLIDE_IDS.index(sid) for sid in excluded_template
         if TEMPLATE_SLIDE_IDS.index(sid) < len(prs.slides)],
        reverse=True,
    )
    for idx in indices_to_delete:
        _delete_slide(prs, idx)

    # --- Reorder remaining template slides ---
    # After deletion, the remaining slides are in template order.
    # We need to reorder to match `included`.
    # Rebuild slide_map after deletion (slides are now re-indexed).
    _reorder_slides(prs, included)

    # --- Insert / append optional slides at correct deck positions ---
    # current_order tracks the live slide sequence as we add optional slides.
    # We use it to find the right insertion index for each new slide.
    current_order = [sid for sid in included if sid in TEMPLATE_SLIDE_IDS]

    # Positioning rules: optional slide → insert BEFORE this template slide
    # (or after the last timeline slide for the scenario variants).
    _INSERT_BEFORE: Dict[str, str] = {
        "acf_composition":      "timeline_chart",     # right before the ACF compliance timeline
        "invalid_vin":          "data_needs",          # before data needs
        "department_summary":   "infra_costs_chart",   # appendix area, before infra
        "facility_summary":     "infra_costs_chart",   # appendix area, before infra
    }

    def _find_insertion_idx(anchor_id: str) -> int:
        """Return the current index at which to insert (before anchor_id)."""
        try:
            return current_order.index(anchor_id)
        except ValueError:
            return len(prs.slides)  # append if anchor not present

    def _find_after_last_timeline() -> int:
        """Insert after the last timeline-related slide (template or optional)."""
        last = -1
        for i, sid in enumerate(current_order):
            if sid == "timeline_chart" or sid.startswith("timeline_"):
                last = i
        return last + 1 if last >= 0 else len(prs.slides)

    for opt_id in (profile.optional_slides or []):
        n_before = len(prs.slides)
        success = _append_optional_slide(prs, opt_id, vehicles, profile, scenario_results)
        if not success or len(prs.slides) <= n_before:
            continue  # slide was not added (conditional slide, no data, etc.)

        # The new slide is at index n_before (last slide)
        if opt_id in _INSERT_BEFORE:
            target = _find_insertion_idx(_INSERT_BEFORE[opt_id])
            _move_slide_to_index(prs, n_before, target)
            current_order.insert(target, opt_id)
        elif opt_id.startswith("timeline_"):
            target = _find_after_last_timeline()
            _move_slide_to_index(prs, n_before, target)
            current_order.insert(target, opt_id)
        else:
            # Append at end (age_analysis, replacement_table, data_quality, etc.)
            current_order.append(opt_id)

    # --- Determine output path ---
    if out_path is None:
        ts = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        client_slug = (profile.client_name or "Fleet").replace(" ", "_")[:30]
        filename = f"{client_slug}_{profile.presentation_type}_{ts}.pptx"
        out_path = str(Path(EXPORT_DIR) / filename)

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    logger.info("[export_presentation] Saved to: %s", out_path)
    return os.path.abspath(out_path)


###############################################################################
# PDF export
###############################################################################

def export_pdf(pptx_path: str) -> Optional[str]:
    """
    Convert a .pptx file to PDF.
    Tries LibreOffice (soffice) first; falls back to macOS open-with-PowerPoint.
    Returns the PDF path on success, or None with a descriptive message.
    """
    pptx_path = os.path.abspath(pptx_path)
    pdf_path = Path(pptx_path).with_suffix(".pdf")
    out_dir = str(pdf_path.parent)

    # 1. Try LibreOffice headless
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode == 0 and pdf_path.exists():
                logger.info("[export_pdf] PDF created via LibreOffice: %s", pdf_path)
                return str(pdf_path)
            logger.warning("[export_pdf] LibreOffice failed: %s", result.stderr)
        except Exception as e:
            logger.warning("[export_pdf] LibreOffice error: %s", e)

    # 2. Try pptx2pdf (wraps LibreOffice/PowerPoint)
    try:
        from pptx2pdf import convert  # type: ignore
        convert(pptx_path, str(pdf_path))
        if pdf_path.exists():
            logger.info("[export_pdf] PDF created via pptx2pdf: %s", pdf_path)
            return str(pdf_path)
    except ImportError:
        pass
    except Exception as e:
        logger.warning("[export_pdf] pptx2pdf error: %s", e)

    # 3. Fallback: open in PowerPoint (macOS)
    import platform
    if platform.system() == "Darwin":
        try:
            subprocess.Popen(["open", "-a", "Microsoft PowerPoint", pptx_path])
        except Exception:
            pass
    return None  # Signal to caller to show fallback dialog


###############################################################################
# Backward-compatibility shim
###############################################################################

def export_prelim_deck(data: dict, template_path: Optional[str] = None,
                       out_path: Optional[str] = None,
                       slide_config=None) -> str:
    """
    Backward-compatible shim → delegates to export_presentation().
    The `data` dict should contain 'fleet' or 'vehicles', plus optional profile keys.
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is required. Run: pip install python-pptx")

    # Extract fleet
    fleet_raw = data.get("fleet") or data.get("vehicles") or []
    if isinstance(fleet_raw, Fleet):
        vehicles = fleet_raw.vehicles
    elif isinstance(fleet_raw, list):
        vehicles = fleet_raw
    else:
        vehicles = []

    # Build a minimal profile from the legacy data dict
    profile = PresentationProfile()
    profile.client_name = data.get("client_name", "")
    profile.presentation_type = data.get("stage", "Kickoff")
    profile.included_slides = list(DEFAULT_SLIDE_IDS)
    profile.template_path = template_path

    return export_presentation(
        fleet_data=vehicles,
        profile=profile,
        out_path=out_path,
        template_path=template_path,
    )


###############################################################################
# Legacy helpers (kept for backward compatibility with analysis_panel.py)
###############################################################################

def _get_template_presentation(template_path: Optional[str] = None) -> Presentation:
    """Load a .pptx/.potx template or create blank (legacy helper)."""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx not available")
    if template_path and os.path.isfile(template_path):
        return Presentation(template_path)
    tp = DEFAULT_TEMPLATE_PATH
    if os.path.isfile(tp):
        return Presentation(tp)
    return Presentation()


def _extract_fleet_data(data: dict) -> dict:
    """Normalise the fleet data dict (legacy helper)."""
    fleet = data.get("fleet") or data.get("vehicles") or []
    if isinstance(fleet, Fleet):
        vehicles = fleet.vehicles
        name = fleet.name
    elif isinstance(fleet, list):
        vehicles = fleet
        name = data.get("fleet_name", "Fleet Analysis")
    else:
        vehicles = []
        name = "Fleet Analysis"
    return {"vehicles": vehicles, "fleet_name": name, **data}
