"""
tests/test_pptx_smoke.py

End-to-end headless smoke test for export_presentation().

Builds a representative test fleet in memory (covering all 4 ACF categories + ZEV),
calls export_presentation(), then uses python-pptx to verify the structural and
visual properties of the output deck.

Run: python -m pytest tests/test_pptx_smoke.py -v
"""

import sys
import os
import pytest

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from settings import TEMPLATE_SLIDE_IDS, DEFAULT_SLIDE_IDS, DEFAULT_TEMPLATE_PATH

try:
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# ---------------------------------------------------------------------------
# Module-level skip conditions
# ---------------------------------------------------------------------------

pytestmark = [
    pytest.mark.skipif(not PPTX_AVAILABLE,
                       reason="python-pptx not installed"),
    pytest.mark.skipif(not os.path.isfile(DEFAULT_TEMPLATE_PATH),
                       reason=f"Template not found: {DEFAULT_TEMPLATE_PATH}"),
]


# ---------------------------------------------------------------------------
# Fleet builder
# ---------------------------------------------------------------------------

def _make_test_fleet():
    """Build a minimal fleet covering ZEV + all 4 ACF categories."""
    from tests.conftest import make_fleet_vehicle
    from data.models import Fleet

    fleet = Fleet(name="Smoke Test Fleet")

    # ZEV — already electric, no EV year needed
    zev = make_fleet_vehicle(
        vin="5YJ3E1EA1LF000001",
        vid_overrides=dict(
            fuel_type="Battery Electric Vehicle (BEV)",
            body_class="Sedan/Saloon",
            gvwr="Class 1: 0 - 6,000 lb",
        ),
        fuel_overrides=dict(combined_mpg=0.0, co2_primary=0.0),
        annual_mileage=10000.0,
    )
    zev.custom_fields["_acf_code"] = "ZEV"
    zev.custom_fields["ACF Category"] = "Already Zero-Emission"
    zev.processing_success = True

    # Cat A — light-duty, exempt
    cat_a = make_fleet_vehicle(
        vin="1HGBH41JXMN000002",
        vid_overrides=dict(
            body_class="Sedan/Saloon",
            gvwr="Class 1: 0 - 6,000 lb",
            fuel_type="Gasoline",
        ),
        fuel_overrides=dict(combined_mpg=28.0, co2_primary=316.0),
        annual_mileage=12000.0,
    )
    cat_a.custom_fields["_acf_code"] = "A"
    cat_a.custom_fields["ACF Category"] = "Light-Duty (Exempt)"
    cat_a.custom_fields["Proposed EV Year"] = "2030"
    cat_a.processing_success = True

    # Cat B #1 — mandate-subject medium duty (drives GHG + ACF timeline charts)
    cat_b1 = make_fleet_vehicle(
        vin="1FDUF5HT4FEB000003",
        vid_overrides=dict(
            body_class="Pickup",
            gvwr="Class 4: 14,001 - 16,000 lb",
            fuel_type="Diesel",
        ),
        fuel_overrides=dict(combined_mpg=12.0, co2_primary=741.0),
        annual_mileage=18000.0,
        odometer=95000.0,
    )
    cat_b1.custom_fields["_acf_code"] = "B"
    cat_b1.custom_fields["ACF Category"] = "ACF Mandate-Subject"
    cat_b1.custom_fields["Proposed EV Year"] = "2035"
    cat_b1.custom_fields["_ev_purchase_price"] = 65000.0
    cat_b1.custom_fields["_ice_purchase_price"] = 45000.0
    cat_b1.processing_success = True

    # Cat B #2 — another mandate vehicle (more chart data points)
    cat_b2 = make_fleet_vehicle(
        vin="1FDUF5HT4FEB000004",
        vid_overrides=dict(
            body_class="Pickup",
            gvwr="Class 5: 16,001 - 19,500 lb",
            fuel_type="Diesel",
        ),
        fuel_overrides=dict(combined_mpg=10.0, co2_primary=890.0),
        annual_mileage=20000.0,
        odometer=120000.0,
    )
    cat_b2.custom_fields["_acf_code"] = "B"
    cat_b2.custom_fields["ACF Category"] = "ACF Mandate-Subject"
    cat_b2.custom_fields["Proposed EV Year"] = "2031"
    cat_b2.custom_fields["_ev_purchase_price"] = 75000.0
    cat_b2.custom_fields["_ice_purchase_price"] = 55000.0
    cat_b2.processing_success = True

    # Cat C — body-type exempt (dump truck)
    cat_c = make_fleet_vehicle(
        vin="3C6UR5DL4JG000005",
        vid_overrides=dict(
            body_class="Dump",
            gvwr="Class 7: 26,001 - 33,000 lb",
            fuel_type="Diesel",
        ),
        fuel_overrides=dict(combined_mpg=8.0, co2_primary=1100.0),
        annual_mileage=10000.0,
    )
    cat_c.custom_fields["_acf_code"] = "C"
    cat_c.custom_fields["ACF Category"] = "Body-Type Exempt"
    cat_c.custom_fields["Proposed EV Year"] = "2038"
    cat_c.processing_success = True

    # Cat D — emergency vehicle (PPV)
    cat_d = make_fleet_vehicle(
        vin="1FM5K8AR6MNA000006",
        vid_overrides=dict(
            body_class="Sport Utility Vehicle (SUV)/Multi-Purpose Vehicle (MPV)",
            gvwr="Class 2: 6,001 - 10,000 lb",
            fuel_type="Gasoline",
            trim="PPV",
        ),
        fuel_overrides=dict(combined_mpg=16.0, co2_primary=556.0),
        annual_mileage=25000.0,
    )
    cat_d.custom_fields["_acf_code"] = "D"
    cat_d.custom_fields["ACF Category"] = "Emergency Vehicle"
    cat_d.custom_fields["Proposed EV Year"] = "2040"
    cat_d.processing_success = True

    fleet.vehicles = [zev, cat_a, cat_b1, cat_b2, cat_c, cat_d]
    return fleet


def _make_test_profile():
    """Build a PresentationProfile with all default template slides selected."""
    from data.models import PresentationProfile
    p = PresentationProfile()
    p.client_name = "City of Testville"
    p.presenter_name = "Jane Analyst"
    p.meeting_date = "March 2026"
    p.included_slides = list(DEFAULT_SLIDE_IDS)   # all 15 template slides
    p.optional_slides = []                          # no optional slides
    return p


# ---------------------------------------------------------------------------
# Slide / chart lookup helpers
# ---------------------------------------------------------------------------

def _slide_for(prs: "Presentation", slide_id: str):
    """Return the slide at the TEMPLATE_SLIDE_IDS position for slide_id.

    Valid after export_presentation() with all default slides included (no
    optional slides appended) because the output order matches TEMPLATE_SLIDE_IDS.
    """
    idx = TEMPLATE_SLIDE_IDS.index(slide_id)
    if idx < len(prs.slides):
        return prs.slides[idx]
    return None


def _charts_in(slide) -> list:
    """Return list of Chart objects for every chart shape in a slide."""
    return [shape.chart for shape in slide.shapes if shape.has_chart]


def _all_text(slide) -> str:
    """Concatenate all visible text in a slide."""
    return " ".join(
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame
    )


# ---------------------------------------------------------------------------
# Test class — fixtures generate the PPTX once for the whole class
# ---------------------------------------------------------------------------

class TestPptxSmoke:
    """Structural and content smoke tests for the generated PPTX."""

    @pytest.fixture(scope="class")
    def output_path(self, tmp_path_factory):
        """Generate the PPTX once and return its file path."""
        from powerpoint_export import export_presentation
        fleet = _make_test_fleet()
        profile = _make_test_profile()
        tmp = tmp_path_factory.mktemp("pptx")
        out = str(tmp / "smoke_test.pptx")
        result = export_presentation(fleet, profile=profile, out_path=out)
        # result is a dict with "path" key (and stats); extract path for fixtures
        if isinstance(result, dict):
            return result["path"]
        return result

    @pytest.fixture(scope="class")
    def prs(self, output_path):
        """Open the generated presentation and return the Presentation object."""
        return Presentation(output_path)

    # ------------------------------------------------------------------
    # File integrity
    # ------------------------------------------------------------------

    def test_file_exists(self, output_path):
        """Generated file must exist on disk."""
        assert os.path.isfile(output_path), f"PPTX not created: {output_path}"

    def test_file_is_valid_pptx(self, output_path):
        """File must open without error (valid zip / OOXML structure)."""
        prs = Presentation(output_path)
        assert prs is not None

    # ------------------------------------------------------------------
    # Slide count
    # ------------------------------------------------------------------

    def test_slide_count_equals_template(self, prs):
        """With all 15 default slides selected, output must have exactly 15 slides."""
        expected = len(TEMPLATE_SLIDE_IDS)
        actual = len(prs.slides)
        assert actual == expected, (
            f"Expected {expected} slides (all default), got {actual}"
        )

    # ------------------------------------------------------------------
    # Cover slide
    # ------------------------------------------------------------------

    def test_cover_slide_contains_client_name(self, prs):
        """Cover slide text must include the client name from the profile."""
        cover = _slide_for(prs, "cover")
        assert cover is not None, "cover slide not found at expected index"
        text = _all_text(cover)
        assert "City of Testville" in text, (
            f"Client name missing from cover slide. Found text: {text!r}"
        )

    # ------------------------------------------------------------------
    # Key Findings
    # ------------------------------------------------------------------

    def test_key_findings_has_bullets(self, prs):
        """Key Findings body placeholder must have ≥ 3 non-empty bullet lines."""
        slide = _slide_for(prs, "key_findings")
        assert slide is not None, "key_findings slide not found"

        # Placeholder idx 1 = content body
        content_ph = next(
            (ph for ph in slide.placeholders if ph.placeholder_format.idx == 1),
            None,
        )
        assert content_ph is not None, "key_findings body placeholder (idx=1) not found"

        non_empty = [
            p.text.strip()
            for p in content_ph.text_frame.paragraphs
            if p.text.strip()
        ]
        assert len(non_empty) >= 3, (
            f"Key Findings has only {len(non_empty)} bullet(s); need ≥ 3.\n"
            f"Bullets found: {non_empty}"
        )

    # ------------------------------------------------------------------
    # ACF Electrification Timeline chart
    # ------------------------------------------------------------------

    def test_acf_timeline_chart_exists(self, prs):
        """timeline_chart slide must contain at least one embedded chart."""
        slide = _slide_for(prs, "timeline_chart")
        assert slide is not None
        charts = _charts_in(slide)
        assert len(charts) >= 1, (
            "No chart shape found in timeline_chart slide. "
            "_add_acf_electrification_chart may have returned False (no EV years set?)."
        )

    def test_acf_timeline_chart_has_four_series(self, prs):
        """ACF timeline chart must always contain exactly 4 series (A/B/C/D)."""
        slide = _slide_for(prs, "timeline_chart")
        charts = _charts_in(slide)
        assert charts, "No chart in timeline_chart slide"
        n = len(list(charts[0].series))
        assert n == 4, (
            f"ACF timeline chart should have 4 series for consistent legend; "
            f"found {n}."
        )

    # ------------------------------------------------------------------
    # GHG Emissions chart
    # ------------------------------------------------------------------

    def test_ghg_chart_exists(self, prs):
        """emissions_chart slide must contain at least one embedded chart."""
        slide = _slide_for(prs, "emissions_chart")
        assert slide is not None
        charts = _charts_in(slide)
        assert len(charts) >= 1, (
            "No chart in emissions_chart slide. "
            "_add_ghg_emissions_chart may have returned False."
        )

    def test_ghg_chart_has_three_series(self, prs):
        """GHG chart must have 3 series when Cat B vehicles are present."""
        slide = _slide_for(prs, "emissions_chart")
        charts = _charts_in(slide)
        assert charts, "No chart in emissions_chart slide"
        n = len(list(charts[0].series))
        assert n == 3, (
            f"GHG chart should have 3 series (Baseline / M·H Duty Only / "
            f"Whole Fleet); found {n}. Are Cat B vehicles missing _acf_code='B'?"
        )

    def test_ghg_chart_value_axis_title(self, prs):
        """GHG value-axis title must contain 'Metric Tons'."""
        slide = _slide_for(prs, "emissions_chart")
        charts = _charts_in(slide)
        assert charts, "No chart in emissions_chart slide"
        chart = charts[0]
        assert chart.value_axis.has_title, "GHG chart value axis missing title"
        title = chart.value_axis.axis_title.text_frame.text
        assert "Metric Tons" in title, (
            f"GHG axis title should contain 'Metric Tons', got {title!r}"
        )

    # ------------------------------------------------------------------
    # TCO chart
    # ------------------------------------------------------------------

    def test_tco_chart_exists(self, prs):
        """tco_chart slide must contain at least one embedded chart."""
        slide = _slide_for(prs, "tco_chart")
        assert slide is not None
        charts = _charts_in(slide)
        assert len(charts) >= 1, (
            "No chart in tco_chart slide. "
            "add_tco_comparison_chart may have returned False "
            "(vehicles missing _ev_purchase_price / combined_mpg?)."
        )

    def test_tco_chart_is_clustered_column(self, prs):
        """TCO chart must be COLUMN_CLUSTERED (not COLUMN_STACKED)."""
        slide = _slide_for(prs, "tco_chart")
        charts = _charts_in(slide)
        assert charts, "No chart in tco_chart slide"
        chart_type = charts[0].chart_type
        assert chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED, (
            f"TCO chart should be COLUMN_CLUSTERED, got {chart_type!r}. "
            f"Was add_tco_comparison_chart() updated in Phase 24?"
        )

    def test_tco_chart_value_axis_title(self, prs):
        """TCO value-axis title must contain 'Cost'."""
        slide = _slide_for(prs, "tco_chart")
        charts = _charts_in(slide)
        assert charts, "No chart in tco_chart slide"
        chart = charts[0]
        assert chart.value_axis.has_title, "TCO chart value axis missing title"
        title = chart.value_axis.axis_title.text_frame.text
        assert "Cost" in title, (
            f"TCO axis title should contain 'Cost', got {title!r}"
        )


# ---------------------------------------------------------------------------
# Helpers for optional-slide tests
# ---------------------------------------------------------------------------

def _slide_by_title_fragment(prs: "Presentation", fragment: str):
    """Return the first slide whose *title placeholder* text contains *fragment*.

    Only checks placeholder idx=0 (title) so body text in other slides cannot
    cause false matches.
    """
    for slide in prs.slides:
        title_ph = next(
            (ph for ph in slide.placeholders if ph.placeholder_format.idx == 0),
            None,
        )
        if title_ph and title_ph.has_text_frame and fragment in title_ph.text_frame.text:
            return slide
    return None


def _make_optional_fleet_with_scenarios():
    """Build test fleet and run compare_scenarios() for optional slide fixtures."""
    from analysis.scenarios import compare_scenarios
    from analysis.electrification_timeline import assign_electrification_years

    fleet = _make_test_fleet()
    assign_electrification_years(fleet.vehicles, end_year=2040)
    scenario_results = compare_scenarios(
        fleet.vehicles,
        scenario_names=["moderate", "aggressive", "conservative"],
    )
    return fleet, scenario_results


# ---------------------------------------------------------------------------
# Optional slide tests
# ---------------------------------------------------------------------------

class TestPptxOptionalSlides:
    """Smoke tests for optional slides appended to the presentation.

    Builds a single PPTX with three optional slides enabled:
      - timeline_milestone  (Cat B GVWR chart)
      - scenario_co2        (CO2 trajectory chart)
      - scenario_investment (Cumulative investment chart)

    The fixture is class-scoped so the PPTX is generated once for all tests.
    """

    @pytest.fixture(scope="class")
    def optional_output(self, tmp_path_factory):
        """Generate a PPTX with optional slides; return (Presentation, base_count)."""
        from powerpoint_export import export_presentation
        from data.models import PresentationProfile

        fleet, scenario_results = _make_optional_fleet_with_scenarios()

        profile = PresentationProfile()
        profile.client_name = "Optional Slides Test"
        profile.included_slides = list(DEFAULT_SLIDE_IDS)
        # Process in this order: scenario_co2 and scenario_investment append at end;
        # timeline_milestone inserts after timeline_chart (middle of deck).
        profile.optional_slides = ["scenario_co2", "scenario_investment", "timeline_milestone"]

        tmp = tmp_path_factory.mktemp("pptx_opt")
        out = str(tmp / "optional_smoke.pptx")
        result = export_presentation(
            fleet,
            profile=profile,
            out_path=out,
            scenario_results=scenario_results,
        )
        path = result["path"] if isinstance(result, dict) else result
        prs = Presentation(path)
        return prs, len(DEFAULT_SLIDE_IDS)

    # ------------------------------------------------------------------
    # Slide count
    # ------------------------------------------------------------------

    def test_optional_slides_increase_count(self, optional_output):
        """Three optional slides must raise the total slide count above the base."""
        prs, base_count = optional_output
        assert len(prs.slides) == base_count + 3, (
            f"Expected {base_count + 3} slides (15 base + 3 optional), "
            f"got {len(prs.slides)}"
        )

    # ------------------------------------------------------------------
    # scenario_co2 optional slide
    # ------------------------------------------------------------------

    def test_scenario_co2_slide_present(self, optional_output):
        """A slide titled 'Annual Fleet Emissions by Scenario' must exist."""
        prs, _ = optional_output
        slide = _slide_by_title_fragment(prs, "Annual Fleet Emissions by Scenario")
        assert slide is not None, (
            "scenario_co2 slide not found — title 'Annual Fleet Emissions by Scenario' "
            "missing from all slides."
        )

    def test_scenario_co2_slide_has_chart(self, optional_output):
        """scenario_co2 slide must contain at least one embedded chart."""
        prs, _ = optional_output
        slide = _slide_by_title_fragment(prs, "Annual Fleet Emissions by Scenario")
        assert slide is not None, "scenario_co2 slide not found"
        charts = _charts_in(slide)
        assert len(charts) >= 1, (
            "scenario_co2 slide exists but has no chart shape. "
            "add_co2_trajectory_chart() may have returned False."
        )

    # ------------------------------------------------------------------
    # scenario_investment optional slide
    # ------------------------------------------------------------------

    def test_scenario_investment_slide_present(self, optional_output):
        """A slide titled 'Cumulative Fleet Investment by Scenario' must exist."""
        prs, _ = optional_output
        slide = _slide_by_title_fragment(prs, "Cumulative Fleet Investment by Scenario")
        assert slide is not None, (
            "scenario_investment slide not found — title fragment "
            "'Cumulative Fleet Investment by Scenario' missing from all slides."
        )

    def test_scenario_investment_slide_has_chart(self, optional_output):
        """scenario_investment slide must contain at least one embedded chart."""
        prs, _ = optional_output
        slide = _slide_by_title_fragment(prs, "Cumulative Fleet Investment by Scenario")
        assert slide is not None, "scenario_investment slide not found"
        charts = _charts_in(slide)
        assert len(charts) >= 1, (
            "scenario_investment slide exists but has no chart shape. "
            "add_cumulative_investment_chart() may have returned False."
        )

    # ------------------------------------------------------------------
    # timeline_milestone optional slide
    # ------------------------------------------------------------------

    def test_timeline_milestone_slide_present(self, optional_output):
        """A slide titled 'ZEV Milestone Option' must exist."""
        prs, _ = optional_output
        slide = _slide_by_title_fragment(prs, "ZEV Milestone Option")
        assert slide is not None, (
            "timeline_milestone slide not found — title fragment 'ZEV Milestone Option' "
            "missing from all slides."
        )

    def test_timeline_milestone_slide_has_chart(self, optional_output):
        """timeline_milestone slide must contain at least one embedded chart."""
        prs, _ = optional_output
        slide = _slide_by_title_fragment(prs, "ZEV Milestone Option")
        assert slide is not None, "timeline_milestone slide not found"
        charts = _charts_in(slide)
        assert len(charts) >= 1, (
            "timeline_milestone slide exists but has no chart shape. "
            "_add_milestone_option_chart() may have returned False (no Cat B vehicles?)."
        )
