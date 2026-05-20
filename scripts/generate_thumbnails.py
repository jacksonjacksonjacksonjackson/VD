"""
scripts/generate_thumbnails.py

Generate PNG slide thumbnail assets for the Present tab card gallery.

Reads each slide's title from assets/template_default.pptx and produces a
320×180 px PNG for every TEMPLATE_SLIDE_ID in settings.py.  The thumbnails
are saved to assets/slide_thumbnails/{slide_id}.png.

Usage:
    python3 scripts/generate_thumbnails.py

Requirements: python-pptx, Pillow (already in requirements.txt)
"""

import sys
import os

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from settings import DEFAULT_TEMPLATE_PATH, TEMPLATE_SLIDE_IDS, ASSETS_DIR

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------

THUMB_W = 320
THUMB_H = 180

# Color palette — dark navy header + light content area
COLOR_HEADER_BG  = (27, 58, 107)     # dark navy
COLOR_HEADER_LINE = (77, 161, 169)   # teal accent stripe
COLOR_CONTENT_BG  = (245, 247, 250)  # off-white
COLOR_TITLE_TEXT  = (255, 255, 255)  # white on header
COLOR_LABEL_TEXT  = (120, 140, 165)  # muted slate (slide ID label)
COLOR_CHART_MOCK  = (210, 220, 235)  # light blue-gray for chart placeholder area
COLOR_BORDER      = (200, 210, 220)  # subtle border

HEADER_H = 46   # height of the dark title bar
STRIPE_H = 3    # height of accent color stripe under header

# Slide IDs that display chart placeholder icons in the content area
CHART_SLIDE_IDS = {
    "timeline_chart", "emissions_chart", "infra_costs_chart", "tco_chart",
}

# ---------------------------------------------------------------------------
# Font helpers — try to load a system font; fall back to default
# ---------------------------------------------------------------------------

def _load_font(size: int, bold: bool = False):
    candidates = [
        "/Library/Fonts/Avenir LT Std Book.otf",
        "/Library/Fonts/Avenir Next.ttc",
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/SFNSDisplay.ttf",
    ]
    for path in candidates:
        if os.path.isfile(path):
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    return ImageFont.load_default()


FONT_TITLE   = _load_font(14, bold=True)
FONT_LABEL   = _load_font(9)


# ---------------------------------------------------------------------------
# Thumbnail renderer
# ---------------------------------------------------------------------------

def _wrap_text(draw, text: str, font, max_width: int) -> list[str]:
    """Wrap text into lines that fit within max_width pixels."""
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = (current + " " + word).strip()
        bbox = draw.textbbox((0, 0), test, font=font)
        if bbox[2] <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [""]


def render_thumbnail(slide_id: str, title: str) -> Image.Image:
    """Render a 320×180 thumbnail for a slide."""
    img = Image.new("RGB", (THUMB_W, THUMB_H), COLOR_CONTENT_BG)
    draw = ImageDraw.Draw(img)

    # --- Header bar ---
    draw.rectangle([(0, 0), (THUMB_W, HEADER_H)], fill=COLOR_HEADER_BG)

    # Accent stripe below header
    draw.rectangle([(0, HEADER_H), (THUMB_W, HEADER_H + STRIPE_H)], fill=COLOR_HEADER_LINE)

    # --- Title text inside header ---
    pad_x = 10
    max_title_w = THUMB_W - pad_x * 2
    lines = _wrap_text(draw, title, FONT_TITLE, max_title_w)

    # Calculate total text block height
    line_h = 16
    total_text_h = len(lines) * line_h
    text_y = (HEADER_H - total_text_h) // 2

    for line in lines[:2]:  # max 2 lines in header
        draw.text((pad_x, text_y), line, font=FONT_TITLE, fill=COLOR_TITLE_TEXT)
        text_y += line_h

    # --- Content area ---
    content_top = HEADER_H + STRIPE_H + 8
    content_h = THUMB_H - content_top - 8

    if slide_id in CHART_SLIDE_IDS:
        # Draw a simple chart placeholder (bars or line suggestion)
        _draw_chart_mock(draw, 10, content_top, THUMB_W - 20, content_h)
    else:
        # Draw bullet point lines
        _draw_bullet_lines(draw, 12, content_top, THUMB_W - 24, content_h)

    # --- Slide ID label (bottom-right) ---
    label = slide_id.replace("_", " ")
    bbox = draw.textbbox((0, 0), label, font=FONT_LABEL)
    label_w = bbox[2] - bbox[0]
    draw.text((THUMB_W - label_w - 6, THUMB_H - 12), label, font=FONT_LABEL,
              fill=COLOR_LABEL_TEXT)

    # Subtle border
    draw.rectangle([(0, 0), (THUMB_W - 1, THUMB_H - 1)], outline=COLOR_BORDER)

    return img


def _draw_chart_mock(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int):
    """Draw a minimal bar chart suggestion."""
    bar_count = 6
    bar_w = max(4, (w - bar_count * 4) // bar_count)
    colors = [COLOR_HEADER_LINE, COLOR_CHART_MOCK, COLOR_HEADER_BG,
              COLOR_CHART_MOCK, COLOR_HEADER_LINE, COLOR_CHART_MOCK]
    heights = [0.8, 0.55, 0.9, 0.65, 0.45, 0.7]

    for i in range(bar_count):
        bx = x + i * (bar_w + 4)
        bh = int(h * heights[i % len(heights)])
        by = y + h - bh
        color = colors[i % len(colors)]
        draw.rectangle([(bx, by), (bx + bar_w, y + h)], fill=color)

    # X-axis line
    draw.line([(x, y + h), (x + w, y + h)], fill=COLOR_LABEL_TEXT, width=1)


def _draw_bullet_lines(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int):
    """Draw placeholder bullet line bars."""
    line_h = 7
    gap = 5
    bullet_r = 2
    num_lines = min(5, h // (line_h + gap))

    for i in range(num_lines):
        ly = y + i * (line_h + gap)
        # Bullet dot
        draw.ellipse([(x, ly + 2), (x + bullet_r * 2, ly + 2 + bullet_r * 2)],
                     fill=COLOR_HEADER_LINE)
        # Line bar (varying widths for realism)
        line_w = int(w * (0.85 - i * 0.08))
        draw.rectangle([(x + 8, ly + 2), (x + 8 + line_w, ly + line_h)],
                       fill=COLOR_CHART_MOCK)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    out_dir = Path(ASSETS_DIR) / "slide_thumbnails"
    out_dir.mkdir(parents=True, exist_ok=True)

    # Load titles from template
    prs = Presentation(DEFAULT_TEMPLATE_PATH)
    titles: dict[str, str] = {}
    for i, slide in enumerate(prs.slides):
        if i >= len(TEMPLATE_SLIDE_IDS):
            break
        sid = TEMPLATE_SLIDE_IDS[i]
        title = ""
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0 and ph.has_text_frame:
                raw = ph.text_frame.text.strip()
                # Clean up vertical-tab line separators (e.g. timeline_chart has \x0b)
                title = raw.replace("\x0b", " — ").split("\n")[0].strip()
                break
        titles[sid] = title or sid.replace("_", " ").title()

    # Generate thumbnails
    generated = []
    for sid in TEMPLATE_SLIDE_IDS:
        title = titles.get(sid, sid.replace("_", " ").title())
        img = render_thumbnail(sid, title)
        out_path = out_dir / f"{sid}.png"
        img.save(str(out_path), "PNG")
        generated.append(out_path.name)
        print(f"  ✓ {out_path.name:30} — {title[:50]}")

    print(f"\nGenerated {len(generated)} thumbnails → {out_dir}")


if __name__ == "__main__":
    main()
