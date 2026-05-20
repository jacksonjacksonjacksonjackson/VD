"""
powerpoint_charts.py

Native PowerPoint chart generation for the Fleet Electrification Analyzer.
Creates editable charts directly in PowerPoint instead of static images.
"""

import logging
from typing import Dict, List, Any, Optional

try:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from data.models import FleetVehicle, Fleet
from settings import PRIMARY_HEX_1, PRIMARY_HEX_3, SECONDARY_HEX_1

logger = logging.getLogger(__name__)

###############################################################################
# Chart style helper
###############################################################################

def _apply_clean_style(chart, legend_pos=XL_LEGEND_POSITION.BOTTOM if PPTX_AVAILABLE else None,
                       colors: Optional[List[str]] = None) -> None:
    """Apply clean consulting-style formatting: no title, bottom legend, light gridlines."""
    if not PPTX_AVAILABLE:
        return
    chart.has_title = False
    if legend_pos is not None:
        chart.has_legend = True
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    try:
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        chart.value_axis.major_gridlines.format.line.width = 6350  # 0.5pt
    except Exception:
        pass

    _DEFAULT_COLORS = [
        PRIMARY_HEX_1,    # charcoal
        PRIMARY_HEX_3,    # reseda green
        SECONDARY_HEX_1,  # deep orange
        "#7B9E87",        # sage
        "#A8C5DA",        # sky blue
        "#D4A373",        # warm sand
    ]
    palette = colors or _DEFAULT_COLORS
    series_list = list(chart.series)
    for i, series in enumerate(series_list):
        hex_c = palette[i % len(palette)]
        r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
        rgb = RGBColor(r, g, b)
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = rgb
        except Exception:
            pass
        try:
            series.format.line.color.rgb = rgb
        except Exception:
            pass


###############################################################################
# Native PowerPoint Chart Functions
###############################################################################

def add_fleet_composition_chart(slide, vehicles: List[FleetVehicle], 
                               left: float = 1, top: float = 2, 
                               width: float = 8, height: float = 4) -> bool:
    """Add native PowerPoint pie chart for fleet composition by body class."""
    try:
        if not vehicles:
            return False
        
        # Count body classes
        body_counts = {}
        for vehicle in vehicles:
            body_class = getattr(vehicle.vehicle_id, 'body_class', 'Unknown') or 'Unknown'
            body_counts[body_class] = body_counts.get(body_class, 0) + 1
        
        # Sort and limit to top categories
        sorted_counts = sorted(body_counts.items(), key=lambda x: x[1], reverse=True)
        if len(sorted_counts) > 8:
            top_counts = sorted_counts[:7]
            other_count = sum(count for _, count in sorted_counts[7:])
            sorted_counts = top_counts + [("Other", other_count)]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = [item[0] for item in sorted_counts]
        chart_data.add_series('Vehicle Count', [item[1] for item in sorted_counts])
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, 
            Inches(left), Inches(top), 
            Inches(width), Inches(height), 
            chart_data
        ).chart
        
        # Styling
        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        # Data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True

        # Apply brand colors per-point
        series = chart.series[0]
        palette = [PRIMARY_HEX_1, PRIMARY_HEX_3, SECONDARY_HEX_1,
                   '#7B9E87', '#A8C5DA', '#D4A373', '#F38BA8', '#4ECDC4']
        for i, point in enumerate(series.points):
            hex_c = palette[i % len(palette)]
            r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(r, g, b)

        return True
        
    except Exception as e:
        logger.error(f"Failed to create fleet composition chart: {e}")
        return False

def add_emissions_timeline_chart(slide, vehicles: List[FleetVehicle],
                               left: float = 1, top: float = 2,
                               width: float = 8, height: float = 4) -> bool:
    """Add CO2 emissions reduction timeline based on actual Proposed EV Year assignments.

    Computes real baseline emissions from vehicle data, then projects annual
    reductions by removing each vehicle's CO2 in the year it's scheduled for
    replacement (per the electrification timeline).
    """
    try:
        if not vehicles:
            return False

        current_year = datetime.now().year

        # Build per-vehicle emissions and replacement year
        vehicle_emissions = {}  # year -> list of annual CO2 (metric tons)
        baseline_emissions = 0.0

        for v in vehicles:
            annual_mileage = v.annual_mileage or 12000
            co2_per_mile = v.fuel_economy.co2_primary or 0
            if co2_per_mile <= 0:
                mpg = v.fuel_economy.combined_mpg or 0
                co2_per_mile = 8900 / mpg if mpg > 0 else 0
            if co2_per_mile <= 0:
                continue

            v_co2 = (co2_per_mile * annual_mileage) / 1000000
            baseline_emissions += v_co2

            # Get proposed replacement year
            ev_year_str = v.custom_fields.get("Proposed EV Year", "")
            try:
                ev_year = int(ev_year_str)
            except (ValueError, TypeError):
                ev_year = None  # N/A, Exempt, or blank

            if ev_year and current_year < ev_year <= current_year + 15:
                vehicle_emissions.setdefault(ev_year, []).append(v_co2)

        if baseline_emissions <= 0:
            return False

        # Build year-by-year emissions projection
        max_year = max(vehicle_emissions.keys()) if vehicle_emissions else current_year + 10
        years = list(range(current_year, max(current_year + 11, max_year + 1)))
        remaining = baseline_emissions
        emissions_data = []
        for year in years:
            removed = sum(vehicle_emissions.get(year, []))
            remaining -= removed
            emissions_data.append(round(max(0, remaining), 1))

        # Create chart
        chart_data = CategoryChartData()
        chart_data.categories = [str(y) for y in years]
        chart_data.add_series('Fleet CO₂ Emissions (MT)', emissions_data)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart

        chart.has_legend = False
        chart.chart_title.text_frame.text = "Fleet CO₂ Emissions Reduction Timeline"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.value_axis.axis_title.text_frame.text = "CO₂ Emissions (Metric Tons)"
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.category_axis.axis_title.text_frame.text = "Year"
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)

        series = chart.series[0]
        series.format.line.color.rgb = RGBColor(*tuple(int(PRIMARY_HEX_1.lstrip('#')[j:j+2], 16) for j in (0, 2, 4)))
        series.format.line.width = Pt(3)

        return True

    except Exception as e:
        logger.error(f"Failed to create emissions timeline chart: {e}")
        return False

def add_emissions_by_weight_class_chart(slide, vehicles: List[FleetVehicle],
                                       left: float = 1, top: float = 2,
                                       width: float = 8, height: float = 4) -> bool:
    """Add pie chart showing CO2 emissions by vehicle weight class."""
    try:
        if not vehicles:
            return False
        
        # Calculate emissions by weight class
        weight_class_emissions = {
            'Light Duty (≤8,500 lbs)': 0,
            'Medium Duty (8,501-19,500 lbs)': 0,
            'Heavy Duty (19,501-33,000 lbs)': 0,
            'Extra Heavy Duty (>33,000 lbs)': 0
        }
        
        for vehicle in vehicles:
            gvwr = getattr(vehicle.vehicle_id, 'gvwr_pounds', 0)
            annual_mileage = getattr(vehicle, 'annual_mileage', 12000)
            co2_per_mile = getattr(vehicle.fuel_economy, 'co2_primary', 400)
            
            if co2_per_mile > 0:
                # Convert to metric tons per year
                vehicle_emissions = (co2_per_mile * annual_mileage) / 1000000
                
                # Classify by weight
                if gvwr <= 8500:
                    weight_class_emissions['Light Duty (≤8,500 lbs)'] += vehicle_emissions
                elif gvwr <= 19500:
                    weight_class_emissions['Medium Duty (8,501-19,500 lbs)'] += vehicle_emissions
                elif gvwr <= 33000:
                    weight_class_emissions['Heavy Duty (19,501-33,000 lbs)'] += vehicle_emissions
                else:
                    weight_class_emissions['Extra Heavy Duty (>33,000 lbs)'] += vehicle_emissions
        
        # Filter out zero emissions classes
        filtered_emissions = {k: round(v, 1) for k, v in weight_class_emissions.items() if v > 0}
        
        if not filtered_emissions:
            return False
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = list(filtered_emissions.keys())
        chart_data.add_series('CO₂ Emissions (MT/year)', list(filtered_emissions.values()))
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart
        
        # Format chart with improved styling
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.font.size = Pt(10)
        chart.chart_title.text_frame.text = "CO₂ Emissions by Vehicle Weight Class"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        
        # Format data labels
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.bold = True
        
        # Apply brand colors to pie slices
        series = chart.series[0]
        colors = [PRIMARY_HEX_1, PRIMARY_HEX_3, SECONDARY_HEX_1, '#4ECDC4']
        for i, point in enumerate(series.points):
            color_hex = colors[i % len(colors)]
            r, g, b = tuple(int(color_hex.lstrip('#')[j:j+2], 16) for j in (0, 2, 4))
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(r, g, b)
        
        return True
        
    except Exception as e:
        logger.error(f"Failed to create emissions by weight class chart: {e}")
        return False

def add_electrification_timeline_by_weight_chart(slide, vehicles: List[FleetVehicle],
                                               left: float = 1, top: float = 2,
                                               width: float = 8, height: float = 4) -> bool:
    """Add stacked bar chart showing electrification timeline by weight class.

    Uses actual Proposed EV Year from vehicle custom_fields instead of
    hardcoded percentage schedules.
    """
    try:
        if not vehicles:
            return False

        current_year = datetime.now().year

        def _weight_class(v):
            gvwr = getattr(v.vehicle_id, 'gvwr_pounds', 0) or 0
            if gvwr <= 8500:
                return 'Light Duty'
            elif gvwr <= 19500:
                return 'Medium Duty'
            else:
                return 'Heavy Duty'

        # Build weight_class -> year -> count from real data
        data = {}  # {weight_class: {year: count}}
        all_years = set()

        for v in vehicles:
            ev_year_str = v.custom_fields.get("Proposed EV Year", "")
            try:
                ev_year = int(ev_year_str)
            except (ValueError, TypeError):
                continue
            if ev_year <= current_year or ev_year > current_year + 20:
                continue

            wc = _weight_class(v)
            data.setdefault(wc, {})
            data[wc][ev_year] = data[wc].get(ev_year, 0) + 1
            all_years.add(ev_year)

        if not all_years:
            return False

        years = sorted(all_years)
        year_labels = [str(y) for y in years]

        chart_data = CategoryChartData()
        chart_data.categories = year_labels

        # Ensure consistent ordering
        for wc in ['Light Duty', 'Medium Duty', 'Heavy Duty']:
            if wc in data:
                counts = [data[wc].get(y, 0) for y in years]
                chart_data.add_series(wc, counts)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.chart_title.text_frame.text = "Electrification Timeline by Weight Class"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.value_axis.axis_title.text_frame.text = "Vehicles Electrified"
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.category_axis.axis_title.text_frame.text = "Year"
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)

        colors = [PRIMARY_HEX_1, PRIMARY_HEX_3, SECONDARY_HEX_1]
        for i, series in enumerate(chart.series):
            color_hex = colors[i % len(colors)]
            r, g, b = tuple(int(color_hex.lstrip('#')[j:j+2], 16) for j in (0, 2, 4))
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(r, g, b)

        return True

    except Exception as e:
        logger.error(f"Failed to create electrification timeline by weight chart: {e}")
        return False

def add_electrification_timeline_by_body_type_chart(slide, vehicles: List[FleetVehicle],
                                                  left: float = 1, top: float = 2,
                                                  width: float = 8, height: float = 4) -> bool:
    """Add stacked bar chart showing electrification timeline by body type.

    Uses actual Proposed EV Year from vehicle custom_fields, grouped by
    the top 5 body classes.
    """
    try:
        if not vehicles:
            return False

        current_year = datetime.now().year

        # Build body_type -> year -> count from real data
        data = {}  # {body_type: {year: count}}
        all_years = set()

        for v in vehicles:
            ev_year_str = v.custom_fields.get("Proposed EV Year", "")
            try:
                ev_year = int(ev_year_str)
            except (ValueError, TypeError):
                continue
            if ev_year <= current_year or ev_year > current_year + 20:
                continue

            body_type = getattr(v.vehicle_id, 'body_class', 'Unknown') or 'Unknown'
            data.setdefault(body_type, {})
            data[body_type][ev_year] = data[body_type].get(ev_year, 0) + 1
            all_years.add(ev_year)

        if not all_years or not data:
            return False

        # Top 5 body types by total scheduled vehicles
        sorted_types = sorted(data.items(), key=lambda x: sum(x[1].values()), reverse=True)
        top_types = sorted_types[:5]
        if len(sorted_types) > 5:
            # Merge remaining into "Other"
            other = {}
            for _, year_counts in sorted_types[5:]:
                for y, c in year_counts.items():
                    other[y] = other.get(y, 0) + c
            if other:
                top_types.append(("Other", other))

        years = sorted(all_years)
        year_labels = [str(y) for y in years]

        chart_data = CategoryChartData()
        chart_data.categories = year_labels

        for body_type, year_counts in top_types:
            counts = [year_counts.get(y, 0) for y in years]
            chart_data.add_series(body_type, counts)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_STACKED,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)
        chart.chart_title.text_frame.text = "Electrification Timeline by Body Type"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.value_axis.axis_title.text_frame.text = "Vehicles Electrified"
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.category_axis.axis_title.text_frame.text = "Year"
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)

        colors = [PRIMARY_HEX_1, PRIMARY_HEX_3, SECONDARY_HEX_1, '#4ECDC4', '#95E1D3', '#F38BA8', '#A8E6CF', '#FFD93D']
        for i, series in enumerate(chart.series):
            color_hex = colors[i % len(colors)]
            r, g, b = tuple(int(color_hex.lstrip('#')[j:j+2], 16) for j in (0, 2, 4))
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(r, g, b)

        return True

    except Exception as e:
        logger.error(f"Failed to create electrification timeline by body type chart: {e}")
        return False

def add_age_distribution_chart(slide, vehicles: List[FleetVehicle],
                              left: float = 1, top: float = 2,
                              width: float = 8, height: float = 4) -> bool:
    """Add column chart showing fleet age distribution."""
    try:
        if not vehicles:
            return False
        
        # Calculate ages and group into bins
        current_year = datetime.now().year
        age_bins = {}
        
        for vehicle in vehicles:
            try:
                year = int(getattr(vehicle.vehicle_id, 'year', 0))
                if year > 0:
                    age = current_year - year
                    age_group = f"{age} years"
                    age_bins[age_group] = age_bins.get(age_group, 0) + 1
            except (ValueError, TypeError):
                continue
        
        if not age_bins:
            return False
        
        # Sort by age
        sorted_ages = sorted(age_bins.items(), key=lambda x: int(x[0].split()[0]))
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = [item[0] for item in sorted_ages]
        chart_data.add_series('Vehicle Count', [item[1] for item in sorted_ages])
        
        # Add chart to slide
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart
        
        # Format chart with improved styling
        chart.has_legend = False
        chart.chart_title.text_frame.text = "Fleet Age Distribution"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True
        
        # Format axes
        chart.value_axis.axis_title.text_frame.text = "Number of Vehicles"
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.category_axis.axis_title.text_frame.text = "Vehicle Age"
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        
        # Format series with brand color
        series = chart.series[0]
        r, g, b = tuple(int(PRIMARY_HEX_3.lstrip('#')[j:j+2], 16) for j in (0, 2, 4))
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(r, g, b)
        
        return True
        
    except Exception as e:
        logger.error(f"Failed to create age distribution chart: {e}")
        return False

###############################################################################
# Financial Charts (Phase 9C/9D)
###############################################################################

def add_tco_comparison_chart(slide, vehicles: List[FleetVehicle],
                            left: float = 1, top: float = 2,
                            width: float = 8, height: float = 4) -> bool:
    """Add clustered column chart comparing fleet-wide ICE vs EV TCO.

    Reads EV equivalent pricing from custom_fields (set by ev_database.py).
    Shows purchase, fuel, and maintenance cost components side-by-side (not stacked).
    """
    try:
        if not vehicles:
            return False

        from analysis.calculations import (
            calculate_annual_fuel_cost, calculate_annual_ev_cost,
            DEFAULT_ICE_MAINTENANCE, DEFAULT_EV_MAINTENANCE,
            DEFAULT_ANNUAL_MILEAGE, DEFAULT_VEHICLE_LIFESPAN
        )

        ice_purchase = 0.0
        ice_fuel = 0.0
        ice_maint = 0.0
        ev_purchase = 0.0
        ev_fuel = 0.0
        ev_maint = 0.0
        counted = 0

        for v in vehicles:
            ev_price = v.custom_fields.get("_ev_purchase_price")
            ice_price = v.custom_fields.get("_ice_purchase_price")
            if not ev_price or not ice_price:
                continue
            if not v.fuel_economy.combined_mpg:
                continue

            counted += 1
            years = DEFAULT_VEHICLE_LIFESPAN
            mileage = v.annual_mileage or DEFAULT_ANNUAL_MILEAGE

            ice_purchase += float(ice_price)
            ev_purchase += float(ev_price)
            ice_fuel += calculate_annual_fuel_cost(v) * years
            ev_fuel += calculate_annual_ev_cost(v) * years
            ice_maint += mileage * DEFAULT_ICE_MAINTENANCE * years
            ev_maint += mileage * DEFAULT_EV_MAINTENANCE * years

        if counted == 0:
            return False

        # Restructure: categories = cost components, series = ICE vs EV
        # COLUMN_CLUSTERED shows 3 side-by-side pairs (clearer than stacked)
        chart_data = CategoryChartData()
        chart_data.categories = ['Purchase', 'Fuel / Energy', 'Maintenance']
        chart_data.add_series('ICE Fleet', [round(ice_purchase), round(ice_fuel), round(ice_maint)])
        chart_data.add_series('EV Fleet',  [round(ev_purchase), round(ev_fuel), round(ev_maint)])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart

        _apply_clean_style(chart, legend_pos=XL_LEGEND_POSITION.BOTTOM,
                           colors=[SECONDARY_HEX_1, PRIMARY_HEX_3])
        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = "Total Cost ($)"
        except Exception:
            pass
        return True

    except Exception as e:
        logger.error(f"Failed to create TCO comparison chart: {e}")
        return False


def add_payback_timeline_chart(slide, vehicles: List[FleetVehicle],
                               left: float = 1, top: float = 2,
                               width: float = 8, height: float = 4) -> bool:
    """Add line chart showing cumulative ICE vs EV fleet cost over time.

    Uses the year-by-year cash flow model to show when the EV fleet
    becomes cheaper than the ICE fleet (crossover point).
    """
    try:
        if not vehicles:
            return False

        from analysis.calculations import (
            calculate_yearly_cash_flows, DEFAULT_VEHICLE_LIFESPAN
        )

        years_count = DEFAULT_VEHICLE_LIFESPAN
        ice_cumulative = [0.0] * (years_count + 1)
        ev_cumulative = [0.0] * (years_count + 1)
        counted = 0

        for v in vehicles:
            ev_price = v.custom_fields.get("_ev_purchase_price")
            ice_price = v.custom_fields.get("_ice_purchase_price")
            if not ev_price or not ice_price or not v.fuel_economy.combined_mpg:
                continue

            counted += 1
            cf = calculate_yearly_cash_flows(
                vehicle=v,
                ev_purchase_price=float(ev_price),
                ice_purchase_price=float(ice_price),
            )

            for row in cf["yearly_flows"]:
                y = row["year"]
                if y <= years_count:
                    ice_cumulative[y] += row["ice_cumulative"]
                    ev_cumulative[y] += row["ev_cumulative"]

        if counted == 0:
            return False

        year_labels = [f"Year {y}" for y in range(years_count + 1)]

        chart_data = CategoryChartData()
        chart_data.categories = year_labels
        chart_data.add_series('ICE Fleet Cost', [round(v) for v in ice_cumulative])
        chart_data.add_series('EV Fleet Cost', [round(v) for v in ev_cumulative])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(10)
        chart.chart_title.text_frame.text = f"Cumulative Fleet Cost — ICE vs EV ({counted} vehicles)"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.value_axis.axis_title.text_frame.text = "Cumulative Cost ($)"
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(11)
        chart.category_axis.axis_title.text_frame.text = ""

        # ICE line
        ice_series = chart.series[0]
        ice_series.format.line.color.rgb = RGBColor(*tuple(int(SECONDARY_HEX_1.lstrip('#')[j:j+2], 16) for j in (0, 2, 4)))
        ice_series.format.line.width = Pt(3)

        # EV line
        ev_series = chart.series[1]
        ev_series.format.line.color.rgb = RGBColor(*tuple(int(PRIMARY_HEX_3.lstrip('#')[j:j+2], 16) for j in (0, 2, 4)))
        ev_series.format.line.width = Pt(3)

        return True

    except Exception as e:
        logger.error(f"Failed to create payback timeline chart: {e}")
        return False


###############################################################################
# Scenario Comparison Charts (Phase 9E)
###############################################################################

def add_scenario_comparison_chart(slide, scenario_results: list,
                                   left: float = 1, top: float = 2,
                                   width: float = 8, height: float = 4,
                                   metric: str = "vehicles") -> bool:
    """Add a multi-line chart comparing electrification scenarios.

    Args:
        slide: PowerPoint slide object
        scenario_results: List of scenario result dicts (from scenarios.run_scenario)
        left/top/width/height: Chart positioning in inches
        metric: Which metric to chart:
            "vehicles" - cumulative vehicles electrified
            "cost" - cumulative investment
            "co2" - cumulative CO₂ reduction (MT/yr)
            "savings" - cumulative annual savings
    """
    try:
        if not scenario_results:
            return False

        # Collect all years across scenarios
        all_years = set()
        for r in scenario_results:
            if metric == "vehicles":
                all_years.update(r.get("cumulative_vehicles", {}).keys())
            elif metric == "cost":
                all_years.update(r.get("cumulative_cost", {}).keys())
            elif metric == "co2":
                all_years.update(r.get("cumulative_co2_reduction", {}).keys())
            elif metric == "savings":
                all_years.update(r.get("cumulative_savings", {}).keys())

        if not all_years:
            return False

        years = sorted(all_years)
        year_labels = [str(y) for y in years]

        chart_data = CategoryChartData()
        chart_data.categories = year_labels

        # Data key mapping
        data_keys = {
            "vehicles": "cumulative_vehicles",
            "cost": "cumulative_cost",
            "co2": "cumulative_co2_reduction",
            "savings": "cumulative_savings",
        }
        data_key = data_keys.get(metric, "cumulative_vehicles")

        for r in scenario_results:
            values = r.get(data_key, {})
            series_data = [values.get(y, 0) for y in years]
            if metric in ("cost", "savings"):
                series_data = [round(v) for v in series_data]
            elif metric == "co2":
                series_data = [round(v, 1) for v in series_data]
            chart_data.add_series(r["name"], series_data)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data
        ).chart

        # Titles
        titles = {
            "vehicles": "Vehicles Electrified by Scenario",
            "cost": "Cumulative Investment by Scenario",
            "co2": "Cumulative CO₂ Reduction by Scenario (MT/yr)",
            "savings": "Cumulative Annual Savings by Scenario",
        }
        y_labels = {
            "vehicles": "Vehicles",
            "cost": "Total Investment ($)",
            "co2": "CO₂ Reduction (MT/yr)",
            "savings": "Annual Savings ($)",
        }

        chart.chart_title.text_frame.text = titles.get(metric, "Scenario Comparison")
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)

        chart.value_axis.axis_title.text_frame.text = y_labels.get(metric, "")
        chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        chart.category_axis.axis_title.text_frame.text = "Year"
        chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)

        # Color each scenario line
        colors = [PRIMARY_HEX_1, PRIMARY_HEX_3, SECONDARY_HEX_1, '#4ECDC4', '#F38BA8']
        for i, series in enumerate(chart.series):
            color_hex = colors[i % len(colors)]
            r, g, b = tuple(int(color_hex.lstrip('#')[j:j+2], 16) for j in (0, 2, 4))
            series.format.line.color.rgb = RGBColor(r, g, b)
            series.format.line.width = Pt(2.5)

        return True

    except Exception as e:
        logger.error(f"Failed to create scenario comparison chart: {e}")
        return False


###############################################################################
# Scope-scenario line charts (Annual CO₂ trajectory + Cumulative Investment)
###############################################################################

def add_co2_trajectory_chart(
    slide, scenario_results_list: list,
    baseline_co2: float = 0.0,
    left: float = 1.0, top: float = 1.8,
    width: float = 10.0, height: float = 5.0,
) -> bool:
    """Add a LINE_MARKERS chart showing annual remaining fleet emissions by scenario.

    Each series is one scope-based scenario (Minimum Compliance, All Excl. Emergency,
    Whole Fleet).  A dashed baseline series shows no-electrification emissions.

    Args:
        scenario_results_list: list of scenario dicts from compare_scenarios()["scenarios"]
        baseline_co2: annual baseline fleet CO₂ in MT CO₂e (flat line across all years)
    """
    try:
        if not PPTX_AVAILABLE or not scenario_results_list:
            return False

        # Collect all years
        all_years: set = set()
        for r in scenario_results_list:
            all_years.update(r.get("cumulative_co2_reduction", {}).keys())
        if not all_years:
            return False
        years = sorted(all_years)
        year_labels = [str(y) for y in years]

        chart_data = CategoryChartData()
        chart_data.categories = year_labels

        # Baseline (flat)
        chart_data.add_series(
            "Baseline (No Electrification)",
            [round(baseline_co2, 1)] * len(years)
        )

        # Per-scenario annual remaining emissions = baseline − cumulative_co2_reduction[yr]
        for r in scenario_results_list:
            if r.get("total_vehicles", 0) == 0:
                continue
            cum = r.get("cumulative_co2_reduction", {})
            values = [round(max(0.0, baseline_co2 - cum.get(yr, 0)), 1) for yr in years]
            chart_data.add_series(r["name"], values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data,
        ).chart

        chart.has_title = True
        chart.chart_title.text_frame.text = "Annual Fleet Emissions by Scenario"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)

        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = "Metric Tons CO₂e"
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except Exception:
            pass

        # Style series: baseline dashed grey, scenarios in brand colors
        colors = ["#888888", PRIMARY_HEX_1, "#E64A19", "#1B5E20", "#7b52ab"]
        for i, series in enumerate(chart.series):
            c = colors[i % len(colors)]
            r_v, g_v, b_v = tuple(int(c.lstrip("#")[j:j+2], 16) for j in (0, 2, 4))
            series.format.line.color.rgb = RGBColor(r_v, g_v, b_v)
            series.format.line.width = Pt(2.0 if i > 0 else 1.5)
            try:
                if i == 0:  # baseline dashed
                    from pptx.oxml.ns import qn
                    from lxml import etree
                    ln = series.format.line._element
                    dash = etree.SubElement(ln, qn("a:prstDash"))
                    dash.set("val", "dash")
            except Exception:
                pass

        return True

    except Exception as e:
        logger.error(f"Failed to create CO₂ trajectory chart: {e}")
        return False


def add_cumulative_investment_chart(
    slide, scenario_results_list: list,
    left: float = 1.0, top: float = 1.8,
    width: float = 10.0, height: float = 5.0,
) -> bool:
    """Add a LINE_MARKERS chart showing cumulative fleet EV investment by scenario.

    Args:
        scenario_results_list: list of scenario dicts from compare_scenarios()["scenarios"]
    """
    try:
        if not PPTX_AVAILABLE or not scenario_results_list:
            return False

        all_years: set = set()
        for r in scenario_results_list:
            all_years.update(r.get("cumulative_cost", {}).keys())
        if not all_years:
            return False
        years = sorted(all_years)
        year_labels = [str(y) for y in years]

        chart_data = CategoryChartData()
        chart_data.categories = year_labels

        for r in scenario_results_list:
            if r.get("total_vehicles", 0) == 0:
                continue
            cum = r.get("cumulative_cost", {})
            # Scale to $M, round to 1 decimal
            values = [round(cum.get(yr, 0) / 1_000_000, 2) for yr in years]
            chart_data.add_series(r["name"], values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(left), Inches(top),
            Inches(width), Inches(height),
            chart_data,
        ).chart

        chart.has_title = True
        chart.chart_title.text_frame.text = "Cumulative Fleet Investment by Scenario"
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(13)
        chart.chart_title.text_frame.paragraphs[0].font.bold = True

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)

        try:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = "$M (Cumulative)"
            chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except Exception:
            pass

        colors = [PRIMARY_HEX_1, "#E64A19", "#1B5E20", "#7b52ab", "#0077b6"]
        for i, series in enumerate(chart.series):
            c = colors[i % len(colors)]
            r_v, g_v, b_v = tuple(int(c.lstrip("#")[j:j+2], 16) for j in (0, 2, 4))
            series.format.line.color.rgb = RGBColor(r_v, g_v, b_v)
            series.format.line.width = Pt(2.0)

        return True

    except Exception as e:
        logger.error(f"Failed to create cumulative investment chart: {e}")
        return False


###############################################################################
# Phase 24: Fleet Composition by ACF Category + Department/Facility Charts
###############################################################################

def add_acf_category_composition_chart(slide, vehicles: List[FleetVehicle],
                                        left: float = 1.5, top: float = 1.8,
                                        width: float = 10.0, height: float = 5.2) -> bool:
    """Pie chart showing fleet composition by ACF compliance category.

    Uses plain-English labels matching consultant/client language.
    This replaces the body-type pie as the primary composition chart.
    """
    try:
        if not vehicles:
            return False

        ACF_LABELS = {
            "ZEV": "Already Zero-Emission (ZEV)",
            "A":   "Light Duty (Exempt)",
            "B":   "Medium or Heavy Duty (Mandate-Subject)",
            "C":   "Body-Type Exempt",
            "D":   "Emergency Vehicles (Excluded)",
            "":    "Not Classified",
        }
        ACF_COLORS = {
            "ZEV": "#4CAF50",
            "A":   "#A8C5DA",
            "B":   PRIMARY_HEX_1,
            "C":   "#7B9E87",
            "D":   SECONDARY_HEX_1,
            "":    "#D0CCD0",
        }
        ACF_ORDER = ["B", "A", "C", "D", "ZEV", ""]

        counts: Dict[str, int] = {}
        for v in vehicles:
            code = v.custom_fields.get("_acf_code", "")
            counts[code] = counts.get(code, 0) + 1

        labels, values, colors = [], [], []
        for code in ACF_ORDER:
            if counts.get(code, 0) > 0:
                labels.append(ACF_LABELS.get(code, code))
                values.append(counts[code])
                colors.append(ACF_COLORS.get(code, "#D0CCD0"))

        if not labels:
            return False

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Vehicles", values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data,
        ).chart

        chart.has_title = False
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        chart.plots[0].data_labels.font.size = Pt(10)
        chart.plots[0].data_labels.font.bold = True

        series = chart.series[0]
        for i, point in enumerate(series.points):
            hx = colors[i % len(colors)].lstrip("#")
            r, g, b = int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(r, g, b)

        return True

    except Exception as e:
        logger.error("Failed to create ACF category composition chart: %s", e)
        return False


def add_department_summary_chart(slide, vehicles: List[FleetVehicle],
                                  left: float = 1.5, top: float = 1.8,
                                  width: float = 10.0, height: float = 5.2) -> bool:
    """Horizontal stacked bar chart: vehicle count by department.

    Only meaningful when department data is present in the CSV.
    Bars are split into Mandate-Subject (Cat B) vs Other to show compliance exposure.
    Returns False if no department data is found.
    """
    try:
        dept_total: Dict[str, int] = {}
        dept_b: Dict[str, int] = {}
        for v in vehicles:
            dept = (v.department or "").strip() or "Unassigned"
            dept_total[dept] = dept_total.get(dept, 0) + 1
            if v.custom_fields.get("_acf_code") == "B":
                dept_b[dept] = dept_b.get(dept, 0) + 1

        real_depts = {d for d in dept_total if d != "Unassigned"}
        if not real_depts:
            return False

        sorted_depts = sorted(dept_total.items(), key=lambda x: x[1], reverse=True)[:15]
        labels = [d for d, _ in sorted_depts]
        b_vals   = [dept_b.get(d, 0) for d in labels]
        other_vals = [dept_total[d] - dept_b.get(d, 0) for d in labels]

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Mandate-Subject (Cat. B)", b_vals)
        chart_data.add_series("Other Vehicles", other_vals)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data,
        ).chart

        _apply_clean_style(chart, legend_pos=XL_LEGEND_POSITION.BOTTOM,
                           colors=[PRIMARY_HEX_1, "#A8C5DA"])
        return True

    except Exception as e:
        logger.error("Failed to create department summary chart: %s", e)
        return False


def add_facility_summary_chart(slide, vehicles: List[FleetVehicle],
                                left: float = 1.5, top: float = 1.8,
                                width: float = 10.0, height: float = 5.2) -> bool:
    """Horizontal stacked bar chart: vehicle count by facility/location.

    Uses vehicle.location (maps from CSV aliases: facility, station, yard, depot, base).
    Returns False if no location data is found.
    """
    try:
        loc_total: Dict[str, int] = {}
        loc_b: Dict[str, int] = {}
        for v in vehicles:
            loc = (getattr(v, "location", "") or "").strip()
            if not loc:
                continue
            loc_total[loc] = loc_total.get(loc, 0) + 1
            if v.custom_fields.get("_acf_code") == "B":
                loc_b[loc] = loc_b.get(loc, 0) + 1

        if not loc_total:
            return False

        sorted_locs = sorted(loc_total.items(), key=lambda x: x[1], reverse=True)[:15]
        labels = [l for l, _ in sorted_locs]
        b_vals     = [loc_b.get(l, 0) for l in labels]
        other_vals = [loc_total[l] - loc_b.get(l, 0) for l in labels]

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Mandate-Subject (Cat. B)", b_vals)
        chart_data.add_series("Other Vehicles", other_vals)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED,
            Inches(left), Inches(top), Inches(width), Inches(height),
            chart_data,
        ).chart

        _apply_clean_style(chart, legend_pos=XL_LEGEND_POSITION.BOTTOM,
                           colors=[PRIMARY_HEX_1, "#A8C5DA"])
        return True

    except Exception as e:
        logger.error("Failed to create facility summary chart: %s", e)
        return False


###############################################################################
# Chart Configuration and Customization
###############################################################################

class SlideConfiguration:
    """Configuration class for customizable slide generation."""
    
    def __init__(self):
        self.available_slides = {
            'cover': {
                'name': 'Cover Slide',
                'description': 'Fleet name, client info, and date',
                'required': True,
                'charts': []
            },
            'fleet_snapshot': {
                'name': 'Fleet Snapshot KPIs',
                'description': '6 key metrics with baseline costs and emissions',
                'required': True,
                'charts': []
            },
            'fleet_composition': {
                'name': 'Fleet Composition',
                'description': 'Vehicle composition by body class',
                'required': False,
                'charts': ['pie_body_class']
            },
            'emissions_timeline': {
                'name': 'Emissions Reduction Timeline',
                'description': 'CO₂ emissions reduction over time',
                'required': False,
                'charts': ['line_emissions_timeline']
            },
            'emissions_by_weight': {
                'name': 'Emissions by Weight Class',
                'description': 'CO₂ emissions breakdown by vehicle weight',
                'required': False,
                'charts': ['pie_emissions_weight']
            },
            'electrification_timeline_weight': {
                'name': 'Electrification Timeline by Weight',
                'description': 'Electrification schedule by weight class',
                'required': False,
                'charts': ['stacked_bar_weight_timeline']
            },
            'electrification_timeline_body': {
                'name': 'Electrification Timeline by Body Type',
                'description': 'Electrification schedule by body type',
                'required': False,
                'charts': ['stacked_bar_body_timeline']
            },
            'age_analysis': {
                'name': 'Fleet Age Analysis',
                'description': 'Age distribution and statistics',
                'required': False,
                'charts': ['column_age_distribution']
            },
            'data_quality': {
                'name': 'Data Quality & Completeness',
                'description': 'Data completeness analysis',
                'required': False,
                'charts': []
            },
            'next_steps': {
                'name': 'Next Steps & Roadmap',
                'description': 'Implementation timeline and recommendations',
                'required': False,
                'charts': []
            },
            'financial_summary': {
                'name': 'Financial Summary',
                'description': 'TCO comparison and payback analysis',
                'required': False,
                'charts': ['stacked_bar_tco', 'line_payback']
            },
            'executive_recommendations': {
                'name': 'Executive Recommendations',
                'description': 'Data-driven replacement recommendations with top priority vehicles',
                'required': False,
                'charts': []
            },
            'replacement_schedule': {
                'name': 'Replacement Schedule',
                'description': 'Top 12-15 priority vehicles with EV equivalents and savings',
                'required': False,
                'charts': []
            },
            'scenario_comparison': {
                'name': 'Scenario Comparison',
                'description': 'Side-by-side comparison of electrification timeline scenarios',
                'required': False,
                'charts': ['line_scenario_vehicles', 'line_scenario_cost']
            }
        }
        
        self.available_charts = {
            'pie_body_class': {
                'name': 'Fleet Composition (Pie)',
                'description': 'Pie chart of vehicles by body class',
                'function': add_fleet_composition_chart
            },
            'line_emissions_timeline': {
                'name': 'Emissions Timeline (Line)',
                'description': 'Line chart showing CO₂ reduction over time',
                'function': add_emissions_timeline_chart
            },
            'pie_emissions_weight': {
                'name': 'Emissions by Weight (Pie)',
                'description': 'Pie chart of emissions by weight class',
                'function': add_emissions_by_weight_class_chart
            },
            'stacked_bar_weight_timeline': {
                'name': 'Electrification by Weight (Stacked Bar)',
                'description': 'Stacked bar chart of electrification timeline by weight',
                'function': add_electrification_timeline_by_weight_chart
            },
            'stacked_bar_body_timeline': {
                'name': 'Electrification by Body Type (Stacked Bar)',
                'description': 'Stacked bar chart of electrification timeline by body type',
                'function': add_electrification_timeline_by_body_type_chart
            },
            'column_age_distribution': {
                'name': 'Age Distribution (Column)',
                'description': 'Column chart showing fleet age distribution',
                'function': add_age_distribution_chart
            },
            'stacked_bar_tco': {
                'name': 'TCO Comparison (Stacked Bar)',
                'description': 'ICE vs EV total cost of ownership breakdown',
                'function': add_tco_comparison_chart
            },
            'line_payback': {
                'name': 'Payback Timeline (Line)',
                'description': 'Cumulative cost curves showing breakeven point',
                'function': add_payback_timeline_chart
            },
            'line_scenario_vehicles': {
                'name': 'Scenario Comparison — Vehicles (Line)',
                'description': 'Multi-scenario cumulative vehicle electrification',
                'function': add_scenario_comparison_chart
            },
            'line_scenario_cost': {
                'name': 'Scenario Comparison — Cost (Line)',
                'description': 'Multi-scenario cumulative investment',
                'function': add_scenario_comparison_chart
            }
        }
        
        # Default configuration
        self.selected_slides = [
            'cover',
            'fleet_snapshot',
            'fleet_composition',
            'financial_summary',
            'emissions_timeline',
            'emissions_by_weight',
            'electrification_timeline_weight',
            'replacement_schedule',
            'age_analysis',
            'next_steps'
        ]
    
    def get_slide_options(self) -> Dict[str, Dict[str, Any]]:
        """Get available slide options for user selection."""
        return self.available_slides
    
    def get_chart_options(self) -> Dict[str, Dict[str, Any]]:
        """Get available chart options for user selection."""
        return self.available_charts
    
    def set_selected_slides(self, slide_ids: List[str]) -> bool:
        """Set which slides to include in the presentation."""
        # Validate slide IDs
        invalid_slides = [sid for sid in slide_ids if sid not in self.available_slides]
        if invalid_slides:
            logger.error(f"Invalid slide IDs: {invalid_slides}")
            return False
        
        # Ensure required slides are included while preserving user's order
        required_slides = [sid for sid, info in self.available_slides.items() if info['required']]
        seen = set()
        final_slides = []
        for sid in slide_ids + required_slides:
            if sid not in seen:
                final_slides.append(sid)
                seen.add(sid)
        
        self.selected_slides = final_slides
        return True
    
    def get_selected_slides(self) -> List[str]:
        """Get the currently selected slides."""
        return self.selected_slides
